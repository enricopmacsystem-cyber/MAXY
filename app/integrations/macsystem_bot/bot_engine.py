import asyncio
import socket

# Forza IPv4 (ereditato dal bot originale)
socket.setdefaulttimeout(30)
_original_getaddrinfo = socket.getaddrinfo
def _getaddrinfo_ipv4(*args, **kwargs):
    return [info for info in _original_getaddrinfo(*args, **kwargs) if info[0] == socket.AF_INET]
socket.getaddrinfo = _getaddrinfo_ipv4

import os
import re
import shutil
import logging
import hashlib
import json
from logging.handlers import RotatingFileHandler
from dotenv import load_dotenv
from datetime import datetime
from pathlib import Path

# Lock istanza unica PRIMA del caricamento modelli ML/ChromaDB (lento, ~20s)
_EARLY_PID_FILE = Path(__file__).parent / "bot_manuali.pid"
_EARLY_LOCK_FILE = Path(__file__).parent / "bot_manuali.lock"
_early_lock_fd = None


def _early_processo_vivo(pid: int) -> bool:
    if pid <= 0:
        return False
    if os.name == "nt":
        import ctypes
        handle = ctypes.windll.kernel32.OpenProcess(0x1000, False, pid)
        if not handle:
            return False
        try:
            code = ctypes.c_ulong()
            if ctypes.windll.kernel32.GetExitCodeProcess(handle, ctypes.byref(code)):
                return int(code.value) == 259
            return False
        finally:
            ctypes.windll.kernel32.CloseHandle(handle)
    try:
        os.kill(pid, 0)
    except OSError:
        return False
    return True


def _blocca_seconda_istanza_all_avvio() -> None:
    """Eseguito subito all'avvio script, prima di import pesanti."""
    global _early_lock_fd
    if __name__ != "__main__":
        return
    base = Path(__file__).parent
    sospeso = base / "servizi_sospesi.flag"
    if sospeso.exists():
        print(
            "Bot sospeso (servizi_sospesi.flag). "
            "Usa riattiva_automatismi.bat per riattivare."
        )
        raise SystemExit(0)
    stop = base / "bot_manuali.stop"
    if stop.exists():
        try:
            stop.unlink()
        except OSError:
            pass
    if _EARLY_PID_FILE.exists():
        try:
            pid = int(_EARLY_PID_FILE.read_text(encoding="utf-8").strip())
            if _early_processo_vivo(pid):
                print(f"Bot già in esecuzione (PID {pid}). Usa ferma_bot.bat.")
                raise SystemExit(1)
            _EARLY_PID_FILE.unlink(missing_ok=True)
        except ValueError:
            _EARLY_PID_FILE.unlink(missing_ok=True)
    for _tentativo in range(2):
        try:
            _early_lock_fd = os.open(
                _EARLY_LOCK_FILE,
                os.O_CREAT | os.O_EXCL | os.O_RDWR,
            )
            break
        except FileExistsError:
            lock_stale = True
            if _EARLY_LOCK_FILE.exists():
                try:
                    lock_pid = int(_EARLY_LOCK_FILE.read_text(encoding="utf-8").strip() or "0")
                    if lock_pid > 0 and _early_processo_vivo(lock_pid):
                        lock_stale = False
                except (ValueError, OSError):
                    pass
            if _EARLY_PID_FILE.exists():
                try:
                    pid = int(_EARLY_PID_FILE.read_text(encoding="utf-8").strip())
                    if _early_processo_vivo(pid):
                        lock_stale = False
                except ValueError:
                    pass
            if lock_stale and _tentativo == 0:
                try:
                    _EARLY_LOCK_FILE.unlink(missing_ok=True)
                except OSError:
                    pass
                continue
            print("Bot già in esecuzione. Usa ferma_bot.bat prima di riavviare.")
            raise SystemExit(1)
    else:
        raise SystemExit(1)
    os.write(_early_lock_fd, str(os.getpid()).encode())
    os.fsync(_early_lock_fd)


# Lock disabilitato quando importato come libreria dall'Hub Maxy (solo per __main__ Telegram).
if __name__ == "__main__":
    _blocca_seconda_istanza_all_avvio()

try:
    import edge_tts
except ImportError:
    edge_tts = None  # type: ignore[assignment]

from anthropic import Anthropic

try:
    from telegram import Update
    from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, filters, ContextTypes
    from telegram.request import HTTPXRequest
    _TELEGRAM_AVAILABLE = True
except ImportError:
    _TELEGRAM_AVAILABLE = False
    Update = None  # type: ignore[misc, assignment]
    ApplicationBuilder = None  # type: ignore[misc, assignment]

    class ContextTypes:  # type: ignore[no-redef]
        DEFAULT_TYPE = object

import chromadb
from chromadb.utils import embedding_functions
import fitz  # PyMuPDF
import io
import zipfile
import requests as req_lib

# ─────────────────────────────────────────────
# CONFIGURAZIONE
# ─────────────────────────────────────────────
load_dotenv(Path(__file__).parent / ".env")

BOT_TOKEN      = os.getenv("BOT_TOKEN")
ANTHROPIC_KEY  = os.getenv("ANTHROPIC_KEY")
ADMIN_CHAT_ID  = int(os.getenv("ADMIN_CHAT_ID", "0"))

# Cartella locale dove sono salvati i PDF dei manuali
MANUALI_DIR = Path(os.getenv("MANUALI_DIR", "./manuali"))

# Cartella dove ChromaDB salva l'indice vettoriale persistente
CHROMA_DIR  = Path(os.getenv("CHROMA_DIR", "./chroma_db"))

# ─────────────────────────────────────────────
# FAMIGLIE DI PRODOTTI ELMO
# Mappa codici modello → parole chiave per trovare manuali di serie
# MT = Manuale Tecnico, MU = Manuale Utente, QG = Quick Guide, FT = Foglio Tecnico
# ─────────────────────────────────────────────
FAMIGLIE_ELMO = {
    # Serie PROXIMA (PRX) — codici tecnici e nomi estesi
    # Keyword "MT_serie-PROXIMA" matcha i nuovi file MT scaricati (es. ELMO_serie-PROXIMA_MT_MT_serie-PROXIMA_2025-06_36522.pdf)
    "PRX": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA"],
    "PRX80": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX80"],
    "PRX128": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX128"],
    "PRX256": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX256"],
    "PRX1024": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX1024"],
    # Alias nomi estesi per PROXIMA
    "PROXIMA": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA"],
    "PROXIMA80": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX80"],
    "PROXIMA128": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX128"],
    "PROXIMA256": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX256"],
    "PROXIMA1024": ["MT_serie-PROXIMA", "serie-PROXIMA", "PROXIMA", "PRX1024"],
    # Serie SUPERIA (SPR)
    "SPR": ["SUPERIA", "SPR"],
    "SPR256": ["SUPERIA", "SPR256"],
    "SPR512": ["SUPERIA", "SPR512"],
    "SPR2040": ["SUPERIA", "SPR2040"],
    "SUPERIA": ["SUPERIA"],
    # Serie PREGIO
    "PREGIO": ["MT_serie-PREGIO", "Pregio-series", "PREGIO"],
    "PREGIO500": ["MT_serie-PREGIO", "Pregio-series", "PREGIO500"],
    "PREGIO2000": ["MT_serie-PREGIO", "Pregio-series", "PREGIO2000"],
    # Serie VILLEGGIO/VIDOMO
    "VIDOMO": ["VILLEGGIO", "VIDOMO"],
    "VIDOMO2K": ["VILLEGGIO-NG-TRX", "VIDOMO2K"],
    "VICOMPACT": ["VILLEGGIO", "VICOMPACT"],
    "VILLEGGIO": ["VILLEGGIO"],
    # Serie HEKLA
    "HEKLA": ["HEKLA"],
    "HLNET": ["HEKLA", "HLNET"],
    "HLNODE": ["HEKLA", "HLNODE"],
    # Serie HERCOLA
    "HERCOLA": ["HERCOLA"],
    # Serie TACORA
    "TACORA": ["TACORA"],
    "TA2000": ["TACORA", "TA2000"],
    # Periferiche
    "KARMA": ["KARMA"],
    "AURA": ["AURA"],
    "ANIMA": ["ANIMA"],
    "NIRVA": ["NIRVA"],
    "MIDAS": ["MIDAS"],
    "GATEWAY2K": ["GATEWAY2K"],
    "GATEWAY": ["GATEWAY2K"],
    # e-Connect
    "ECONNECT": ["e-Connect", "E-CONNECT"],
    "CONNECT": ["e-Connect", "E-CONNECT"],
}

# Sottostringhe nel nome file PDF per ogni famiglia ELMO (filtro + priorità MT/MP/MU)
ELMO_FILENAME_MARKERS = {
    "PRX": "proxima", "PROXIMA": "proxima", "PRX80": "proxima", "PRX128": "proxima",
    "PRX256": "proxima", "PRX1024": "proxima",
    "PROXIMA80": "proxima", "PROXIMA128": "proxima", "PROXIMA256": "proxima", "PROXIMA1024": "proxima",
    "SPR": "superia", "SUPERIA": "superia", "SPR256": "superia", "SPR512": "superia", "SPR2040": "superia",
    "PREGIO": "pregio", "PREGIO500": "pregio500", "PREGIO2000": "pregio2000",
    "VIDOMO": "vidomo", "VIDOMO2K": "vidomo2k", "VILLEGGIO": "villeggio", "VICOMPACT": "villeggio",
    "HEKLA": "hekla", "HLNET": "hekla", "HLNODE": "hekla",
    "HERCOLA": "hercola", "TACORA": "tacora", "TA2000": "tacora",
    "KARMA": "karma", "ANIMA": "anima", "AURA": "aura", "NIRVA": "nirva", "MIDAS": "midas",
    "GATEWAY2K": "gateway2k", "GATEWAY": "gateway2k",
    "ECONNECT": "e-connect", "CONNECT": "e-connect",
}

# Alias nomi colloquiali → codici tecnici ELMO
ALIAS_ELMO = {
    "PROXIMA": "PRX",
    "PROXIMA80": "PRX80",
    "PROXIMA128": "PRX128",
    "PROXIMA256": "PRX256",
    "PROXIMA1024": "PRX1024",
    "SUPERIA": "SPR",
    "VILLEGGIO": "VIDOMO",
    "TACORA": "TA2000",
}

# Famiglie RISCO — mappa nomi colloquiali a keyword nei nomi file
FAMIGLIE_RISCO = {
    "LIGHTSYS": ["LightSYS", "lightsys"],
    "LIGHTSYS+": ["LightSYS_Plus", "LightSYS+", "LightSYS-Plus", "lightsys-Plus"],
    "LIGHTSYSPLUS": ["LightSYS_Plus", "LightSYS+", "LightSYS-Plus", "lightsys-Plus"],
    "LIGHTSYS2": ["LightSYS2", "LightSYS_2"],
    "LIGHTSYSAIR": ["LightSYS_Air"],
    "WICOMMPRO": ["WiComm_Pro", "wicomm"],
    "WICOMM": ["WiComm_Pro", "wicomm"],
    "AGILITY4": ["Agility_4", "Agility4"],
    "AGILITY": ["Agility_4", "Agility4"],
}

# Alias nomi colloquiali RISCO → chiavi FAMIGLIE_RISCO
ALIAS_RISCO = {
    "LIGHTSYS": "LIGHTSYS",
    "LIGHTSYSPLUS": "LIGHTSYS+",
    "LIGHTSYS2": "LIGHTSYS2",
    "LIGHTSYSAIR": "LIGHTSYSAIR",
    "WICOMM": "WICOMMPRO",
    "AGILITY": "AGILITY4",
}

VOICE = "it-IT-GiuseppeNeural"

# Quanti chunk recuperare per ogni domanda (più è alto = più contesto, ma più token)
TOP_K = 5

# Dimensione dei chunk in caratteri (circa 400 parole)
CHUNK_SIZE  = 1500
CHUNK_OVERLAP = 200

# Parole italiane che il regex può scambiare per codici modello (es. "devo usare" → USARE)
PAROLE_NON_CODICI = {
    "USARE", "DEVO", "DEVI", "STAFFA", "STAFFE", "PALO", "PARA", "PER",
    "INSTALLARE", "INSTALLAZIONE", "MONTARE", "MONTAGGIO",
    "SERVIZI", "SERVIZIO", "CLIENTE", "FINALE", "FORNITORE", "PIATTAFORMA",
    "APPLICAZIONE", "MOBILE", "WEB", "GUIDA", "GUIDE",
    "ESCLUDE", "ESCLUDERE", "ESCLUSIONE", "ATTUALMENTE", "ATTUALE",
}

# Keyword operative per dominio (linguaggio naturale, non dimensionamento Dahua)
KW_OPERATIVE_ANTINTRUSIONE = [
    "esclud", "esclusion", "inser", "disinser", "partizion", "zona", "zone",
    "programm", "utente", "codice", "tastiera", "arm", "disarm", "bypass",
    "centrale", "impianto", "allarme", "rivelatore", "sensore", "sirena",
    "settore", "settori", "inserimento parziale", "browserone", "area ",
]
KW_OPERATIVE_ANTINCENDIO = [
    "pulsante", "rottura vetro", "call point", "rivelatore", "loop ",
    "centrale notifier", "indirizzat", "convenzional", "fumo", "termico",
    "calore", "incendio", "antincendio", "notifier", "esser", "iq8",
    "flexes", "flex es", "esserbus", "modbus", "esib", "bms", "scada",
]
KW_OPERATIVE_VIDEO = [
    "staffa", "staffaggio", "montaggio", "bracket", "junction box",
    "accessori", "come si monta", "come installo", "supporto", "fissaggio",
    "anpr", "targa", "targhe", "varco", "riconoscimento targa", "webhook",
    "http", "cgi", "api", "sdk", "configur", "smartpss", "deterrenza",
    "tioc", "colorvu", "wizsense",
]
KW_OPERATIVE_TUTTE = (
    KW_OPERATIVE_ANTINTRUSIONE + KW_OPERATIVE_ANTINCENDIO + KW_OPERATIVE_VIDEO
)

# Parole che indicano una marca nel testo della domanda
MARCA_MENTION_KEYWORDS = {
    "RISCO": ["risco", "lightsys", "light sys", "agility", "wicomm"],
    "ELMO": ["elmo", "el.mo", "proxima", "superia", "pregio", "vidomo",
             "villeggio", "karma", "anima", "gateway2k", "gateway 2k", "hekla",
             "hercola", "tacora", "nirva", "aura", "browserone", "passlight"],
    "NOTIFIER": ["notifier", "nfs2", "nfs-2", "am-8200", "am8200", "am-8100",
                 "w5a", "m5a", "m3a", "id3000", "vesda", "faast"],
    "ESSER": ["esser", "iq8", "flexes", "flex es", "esserbus", "winmag", "flex control"],
    "DAHUA": ["dahua", "ipc-", "ipc ", "nvr", "hdw", "hfw", "hdb", "xvr", "dhi-",
              "wizsense", "tioc", "acupick"],
}

# Fallback prodotto/manuale se la domanda è operativa ma senza modello esplicito
MARCA_FALLBACK_PATTERN = {
    "RISCO": "LightSYS-2-Full-Installation-Manual-IT",
    "ELMO": "MT_serie-PROXIMA_2025-06_36522",
    "NOTIFIER": "NFS2",
    "ESSER": "FlexES",
}

# (keywords domanda, pattern file, marca, prodotti opzionali)
# Ordine: voci più specifiche prima. prodotti = match anche solo per chiave prodotto estratta.
DOCUMENTO_MANUALE_MAP = [
    # ── RISCO ──
    (["lightsys 2", "lightsys2", "light sys 2"], "LightSYS-2-Full-Installation-Manual-IT", "RISCO", ["LIGHTSYS2"]),
    (["lightsys+", "lightsys plus", "lightsysplus", "light sys plus", "centrale lightsys+"],
     "LightSYS-Plus-Installer-Manual-IT", "RISCO", ["LIGHTSYS+", "LIGHTSYSPLUS"]),
    (["apprend", "sensore", "contatto", "magnetico", "rivelatore", "wireless"],
     "LightSYS-Plus-Installer-Manual-IT", "RISCO", ["LIGHTSYS+", "LIGHTSYSPLUS"]),
    (["lightsys air", "lightsysair"], "LightSYS_Air", "RISCO", ["LIGHTSYSAIR"]),
    (["agility 4", "agility4"], "Agility_4", "RISCO", ["AGILITY4"]),
    (["agility"], "Agility", "RISCO", ["AGILITY", "AGILITY4"]),
    (["wicomm pro", "wicomm"], "WiComm", "RISCO", ["WICOMMPRO", "WICOMM"]),
    (["lightsys", "light sys", "risco"], "LightSYS-2-Full-Installation-Manual-IT", "RISCO", ["LIGHTSYS"]),
    # ── ELMO ──
    (["programm proxima", "configuraz proxima", "mp proxima"], "MP_serie-PROXIMA_2025-04_36523", "ELMO", None),
    (["proxima", "prx80", "prx128", "prx256", "prx1024", "prx"], "MT_serie-PROXIMA_2025-06_36522", "ELMO",
     ["PRX", "PRX80", "PRX128", "PRX256", "PRX1024", "PROXIMA"]),
    (["superia", "spr256", "spr512", "spr2040", "spr"], "SUPERIA-series", "ELMO", ["SPR", "SPR256", "SPR512", "SUPERIA"]),
    (["pregio500", "pregio 500"], "PREGIO500", "ELMO", ["PREGIO500"]),
    (["pregio2000", "pregio 2000"], "PREGIO2000", "ELMO", ["PREGIO2000"]),
    (["pregio"], "MT_serie-PREGIO", "ELMO", ["PREGIO"]),
    (["vidomo2k", "vidomo 2k"], "VIDOMO2K", "ELMO", ["VIDOMO2K"]),
    (["vidomo", "villeggio", "vicompact"], "VIDOMOBTRX_MT", "ELMO", ["VIDOMO", "VILLEGGIO", "VICOMPACT"]),
    (["karma"], "KARMA_MT", "ELMO", ["KARMA"]),
    (["hekla", "hlnet"], "Serie-HEKLA", "ELMO", ["HEKLA", "HLNET", "HLNODE"]),
    (["gateway 2k", "gateway2k", "gateway 2 k"], "MT_GATEWAY2K", "ELMO", ["GATEWAY2K", "GATEWAY"]),
    (["sensori radio", "sensore radio", "ng-trx", "ng trx"],
     "MT_GATEWAY2K", "ELMO", ["GATEWAY2K", "GATEWAY"]),
    (["hekla", "hlnet", "hlnote"], "HEKLA", "ELMO", ["HEKLA", "HLNET", "HLNODE"]),
    (["hercola"], "HERCOLA", "ELMO", ["HERCOLA"]),
    (["tacora", "ta2000"], "TACORA", "ELMO", ["TACORA", "TA2000"]),
    (["anima"], "ANIMA", "ELMO", ["ANIMA"]),
    (["nirva"], "NIRVA", "ELMO", ["NIRVA"]),
    (["aura"], "AURA", "ELMO", ["AURA"]),
    (["e-connect", "econnect", "economico"], "e-Connect", "ELMO", ["ECONNECT", "CONNECT"]),
    (["settore", "settori", "inserimento parziale", "browserone", "elmo"], "MT_serie-PROXIMA_2025-06_36522", "ELMO", None),
    # ── NOTIFIER ──
    (["am-8200n", "am 8200n", "am8200n", "modbus", "mod-bus", "esib", "bms", "scada",
      "supervisione notifier", "seriale notifier"], "AM-8200N manu-prog", "NOTIFIER", None),
    (["am-8200g", "am 8200g", "am8200g"], "AM-8200G manuale di prog", "NOTIFIER", None),
    (["am-8100", "am 8100", "am8100"], "AM-8100 manu-prog", "NOTIFIER", None),
    (["am82-cl", "am82cl"], "AM82-CL manu", "NOTIFIER", None),
    (["am2000", "am-2000", "am 2000"], "AM2000CL-dep", "NOTIFIER", None),
    (["am6000", "am-6000", "am 6000"], "AM6000CL-dep", "NOTIFIER", None),
    (["w5a", "m5a", "m3a", "pulsante", "rottura vetro", "call point",
      "cella frigo", "cella frigorifera"], "M5A-W5A", "NOTIFIER", None),
    (["fsp", "fumo ottico", "rivelatore fumo"], "FSP", "NOTIFIER", None),
    (["fdm", "termico", "rivelatore termico"], "FDM", "NOTIFIER", None),
    (["nfs2", "nfs-2", "centrale notifier", "loop notifier"], "NFS2", "NOTIFIER", None),
    (["id3000", "id 3000", "convenzionale notifier"], "ID3000", "NOTIFIER", None),
    # ── ESSER ──
    (["flex", "flexes", "flex es", "centrale flex", "flex control", "fx2", "fx10", "fx18", "essernet"],
     "FlexES", "ESSER", None),
    (["iq8control", "iq8", "esserbus", "8000c", "9200", "tool8000", "rivelatore esser", "iq8quad", "serie 8000"],
     "esserbus", "ESSER", None),
    (["winmag", "win mag"], "Winmag", "ESSER", None),
    # ── DAHUA ──
    (["staffa", "staffaggio", "montaggio", "bracket", "junction box", "accessori",
      "pfb", "pfa", "pfh", "installare a parete", "come si monta", "come installo",
      "supporto", "fissaggio", "camera accessories", "accessories selection"],
     "Camera-Accessories-Selection", "DAHUA", None),
    (["anpr", "targa", "targhe", "varco", "riconoscimento targa", "lettura targa",
      "white list", "blacklist", "lista targhe", "accesso veicoli"],
     "ANPR", "DAHUA", None),
    (["http", "cgi", "api", "sdk", "webhook", "integrazione", "comando http",
      "notifica http", "action url", "event url", "configurazione telecamera",
      "configurazione camera", "web 5.0", "web 3.0"],
     "Network_Camera_Web_5.0_Operation", "DAHUA", None),
]

# Pattern (testo minuscolo, chiave prodotto) per estrazione da linguaggio naturale
PATTERN_PRODOTTI_COLLOQUIALI = [
    # RISCO
    ("lightsys 2", "LIGHTSYS2"), ("lightsys2", "LIGHTSYS2"),
    ("lightsys+", "LIGHTSYS+"), ("lightsys plus", "LIGHTSYS+"), ("lightsysplus", "LIGHTSYS+"),
    ("lightsys air", "LIGHTSYSAIR"), ("lightsysair", "LIGHTSYSAIR"),
    ("wicomm pro", "WICOMMPRO"), ("wicomm", "WICOMMPRO"),
    ("agility 4", "AGILITY4"), ("agility4", "AGILITY4"), ("agility", "AGILITY4"),
    ("lightsys", "LIGHTSYS"), ("light sys", "LIGHTSYS"),
    # ELMO
    ("proxima 1024", "PRX1024"), ("proxima 256", "PRX256"), ("proxima 128", "PRX128"),
    ("proxima 80", "PRX80"), ("proxima", "PROXIMA"), ("prx1024", "PRX1024"),
    ("prx256", "PRX256"), ("prx128", "PRX128"), ("prx80", "PRX80"), ("prx", "PRX"),
    ("superia", "SUPERIA"), ("spr512", "SPR512"), ("spr256", "SPR256"), ("spr", "SPR"),
    ("pregio 2000", "PREGIO2000"), ("pregio2000", "PREGIO2000"),
    ("pregio 500", "PREGIO500"), ("pregio500", "PREGIO500"), ("pregio", "PREGIO"),
    ("vidomo 2k", "VIDOMO2K"), ("vidomo2k", "VIDOMO2K"),
    ("villeggio", "VIDOMO"), ("vidomo", "VIDOMO"), ("vicompact", "VICOMPACT"),
    ("gateway 2k", "GATEWAY2K"), ("gateway2k", "GATEWAY2K"),
    ("hekla", "HEKLA"), ("hercola", "HERCOLA"), ("tacora", "TACORA"), ("ta2000", "TA2000"),
    ("karma", "KARMA"), ("anima", "ANIMA"), ("nirva", "NIRVA"), ("aura", "AURA"),
    ("e-connect", "ECONNECT"), ("econnect", "ECONNECT"),
    # NOTIFIER (codici corti)
    ("am-8200n", "AM8200N"), ("am8200n", "AM8200N"), ("am-8200g", "AM8200G"),
    ("am-8100", "AM8100"), ("nfs2", "NFS2"), ("nfs-2", "NFS2"), ("id3000", "ID3000"),
    # ESSER
    ("iq8control", "IQ8"), ("flexes", "FLEXES"), ("flex es", "FLEXES"), ("winmag", "WINMAG"),
]

PID_FILE = Path(__file__).parent / "bot_manuali.pid"
STOP_FILE = Path(__file__).parent / "bot_manuali.stop"
LOCK_FILE = Path(__file__).parent / "bot_manuali.lock"
_instance_lock_fd = None

# Cache locale hash file già indicizzati (evita query ChromaDB per ogni file)
HASH_CACHE_FILE = Path("./manuali_hash_cache.json")

def carica_hash_cache() -> dict:
    """Carica la cache degli hash dal file JSON."""
    if HASH_CACHE_FILE.exists():
        try:
            with open(HASH_CACHE_FILE, "r") as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def salva_hash_cache(cache: dict):
    """Salva la cache degli hash nel file JSON."""
    with open(HASH_CACHE_FILE, "w") as f:
        json.dump(cache, f, ensure_ascii=False, indent=2)

# Memoria conversazione per utente (ultimi 6 messaggi)
CONVERSAZIONI = {}
MAX_MESSAGGI  = 6

# ─────────────────────────────────────────────
# API DAHUA PER DOWNLOAD AL VOLO
# ─────────────────────────────────────────────
DAHUA_API_BASE    = "https://www.dahuasecurity.com/frontInterface"
DAHUA_API_SEARCH  = f"{DAHUA_API_BASE}/download_center/document/lan/en"
DAHUA_API_FILZIP  = f"{DAHUA_API_BASE}/download_center/file_zip/lan/en"
DAHUA_HEADERS     = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
    "Referer":    "https://www.dahuasecurity.com/en/support/downloadCenter/documents",
    "Accept":     "application/json, text/plain, */*",
}

if not ANTHROPIC_KEY and __name__ == "__main__":
    raise ValueError("❌ ANTHROPIC_KEY mancante! Controlla il file .env")

# ─────────────────────────────────────────────
# LOGGING (solo logger del bot — evita flood httpx/huggingface nel file)
# ─────────────────────────────────────────────
def _rotate_log_windows_safe(source: str, dest: str):
    """Rotazione log compatibile con Windows (il file può restare aperto)."""
    try:
        os.rename(source, dest)
    except OSError:
        shutil.copy2(source, dest)
        with open(source, "w", encoding="utf-8"):
            pass


_log_fmt = logging.Formatter("%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
logger.propagate = False
logger.handlers.clear()
_console = logging.StreamHandler()
_console.setFormatter(_log_fmt)
logger.addHandler(_console)
_file_handler = RotatingFileHandler(
    "bot_manuali.log",
    maxBytes=50 * 1024 * 1024,
    backupCount=3,
    encoding="utf-8",
    delay=True,
)
_file_handler.rotator = _rotate_log_windows_safe
_file_handler.setFormatter(_log_fmt)
logger.addHandler(_file_handler)

for _noisy in ("httpx", "httpcore", "huggingface_hub", "sentence_transformers", "chromadb"):
    logging.getLogger(_noisy).setLevel(logging.WARNING)

# ─────────────────────────────────────────────
# CHROMADB — DATABASE VETTORIALE
# ─────────────────────────────────────────────
CHROMA_DIR.mkdir(parents=True, exist_ok=True)

chroma_client = chromadb.PersistentClient(path=str(CHROMA_DIR))

# Usa le embedding di sentence-transformers (multilingue, gira in locale, gratis)
embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name="paraphrase-multilingual-MiniLM-L12-v2"
)

collection = chroma_client.get_or_create_collection(
    name="manuali",
    embedding_function=embedding_fn,
    metadata={"hnsw:space": "cosine"},
)

# ─────────────────────────────────────────────
# ESTRAZIONE TESTO DAI PDF
# ─────────────────────────────────────────────
def estrai_testo_pdf(pdf_path: Path) -> str:
    """Estrae testo da PDF usando PyMuPDF, con fallback OCR per PDF scansionati."""
    testo = ""
    try:
        doc = fitz.open(str(pdf_path))
        for page in doc:
            testo += page.get_text()
        doc.close()
    except Exception as e:
        logger.error(f"Errore lettura {pdf_path.name}: {e}")
        return ""

    # Se il testo estratto è troppo poco, prova con OCR
    if len(testo.strip()) < 50:
        try:
            import pytesseract
            from PIL import Image
            import io
            logger.info(f"  OCR in corso per {pdf_path.name}...")
            doc = fitz.open(str(pdf_path))
            testo_ocr = ""
            for page in doc:
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                testo_ocr += pytesseract.image_to_string(img, lang="ita+eng") + "\n"
            doc.close()
            if len(testo_ocr.strip()) > len(testo.strip()):
                logger.info(f"  OCR riuscito: {len(testo_ocr)} caratteri")
                return testo_ocr
        except Exception as e:
            logger.debug(f"OCR fallito per {pdf_path.name}: {e}")

    return testo

def estrai_testo_file(path: Path) -> str:
    """Estrae testo da PDF, TXT, PPTX o DOCX."""
    suffix = path.suffix.lower()
    if suffix == ".txt":
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception as e:
            logger.error(f"Errore lettura {path.name}: {e}")
            return ""
    if suffix in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            testi = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        testi.append(shape.text.strip())
            return "\n".join(testi)
        except Exception as e:
            logger.error(f"Errore lettura PPTX {path.name}: {e}")
            return ""
    if suffix == ".docx":
        try:
            import docx as _docx
            doc = _docx.Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.error(f"Errore lettura DOCX {path.name}: {e}")
            return ""
    if suffix == ".doc":
        # Formato binario Word — non supportato senza antiword/LibreOffice
        logger.warning(f"File .doc non supportato (formato binario): {path.name}")
        return ""
    return estrai_testo_pdf(path)

def chunking(testo: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Divide il testo in chunk sovrapposti per non perdere contesto ai bordi."""
    chunks = []
    start = 0
    while start < len(testo):
        end = start + chunk_size
        chunks.append(testo[start:end])
        start += chunk_size - overlap
    return [c.strip() for c in chunks if len(c.strip()) > 100]

def hash_file(path: Path) -> str:
    """Hash MD5 del file per rilevare modifiche."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            h.update(block)
    return h.hexdigest()

# ─────────────────────────────────────────────
# INDICIZZAZIONE
# ─────────────────────────────────────────────
def indicizza_manuali(forzato: bool = False) -> dict:
    """
    Legge tutti i PDF/TXT in MANUALI_DIR e li indicizza in ChromaDB.
    Usa una cache JSON locale per saltare istantaneamente i file già indicizzati.
    Ritorna un dict con statistiche.
    """
    MANUALI_DIR.mkdir(parents=True, exist_ok=True)
    pdf_files = (list(MANUALI_DIR.glob("**/*.pdf")) +
                 list(MANUALI_DIR.glob("**/*.txt")) +
                 list(MANUALI_DIR.glob("**/*.pptx")) +
                 list(MANUALI_DIR.glob("**/*.PPTX")) +
                 list(MANUALI_DIR.glob("**/*.docx")) +
                 list(MANUALI_DIR.glob("**/*.DOCX")))

    if not pdf_files:
        return {"totale": 0, "nuovi": 0, "saltati": 0, "errori": 0}

    stats = {"totale": len(pdf_files), "nuovi": 0, "saltati": 0, "errori": 0}

    # Carica cache hash — se forzato la ignora
    hash_cache = {} if forzato else carica_hash_cache()

    for pdf_path in pdf_files:
        file_hash = hash_file(pdf_path)
        try:
            cache_key = str(pdf_path.relative_to(MANUALI_DIR))
        except ValueError:
            cache_key = pdf_path.name
        doc_id_prefix = f"{pdf_path.stem}_{file_hash}"

        # Controllo istantaneo sulla cache locale — senza query ChromaDB
        if not forzato and hash_cache.get(cache_key) == file_hash:
            stats["saltati"] += 1
            continue

        # Se il file è cambiato, rimuovi i vecchi chunk da ChromaDB
        if not forzato and cache_key in hash_cache:
            try:
                old_ids = collection.get(where={"source": pdf_path.name})["ids"]
                if old_ids:
                    collection.delete(ids=old_ids)
                    logger.info(f"🔄 {pdf_path.name} aggiornato, reindicizzazione...")
            except Exception:
                pass

        logger.info(f"📖 Indicizzazione: {pdf_path.name}")
        testo = estrai_testo_file(pdf_path)

        if not testo.strip():
            logger.warning(f"⚠️  {pdf_path.name} vuoto o non leggibile.")
            stats["errori"] += 1
            continue

        chunks = chunking(testo)
        logger.info(f"   → {len(chunks)} chunk da {len(testo)} caratteri")

        ids       = [f"{doc_id_prefix}_chunk{i}" for i in range(len(chunks))]
        metadatas = [
            {
                "source":    pdf_path.name,
                "file_hash": file_hash,
                "chunk_idx": i,
                "pagina":    f"~{(i * CHUNK_SIZE) // 2000 + 1}",
            }
            for i in range(len(chunks))
        ]

        try:
            batch = 500
            for b in range(0, len(chunks), batch):
                collection.upsert(
                    ids=ids[b:b+batch],
                    documents=chunks[b:b+batch],
                    metadatas=metadatas[b:b+batch],
                )
            # Aggiorna cache locale dopo successo
            hash_cache[cache_key] = file_hash
            salva_hash_cache(hash_cache)
            stats["nuovi"] += 1
        except Exception as e:
            logger.error(f"Errore ChromaDB per {pdf_path.name}: {e}")
            stats["errori"] += 1

    logger.info(f"✅ Indicizzazione completata: {stats}")
    return stats

# ─────────────────────────────────────────────
# RICERCA E RISPOSTA
# ─────────────────────────────────────────────
def trova_modelli_simili(prefisso: str, max_risultati: int = 10) -> list[str]:
    """
    Cerca nella cartella manuali tutti i modelli che contengono il prefisso.
    Ritorna lista di codici modello unici trovati (es. IPC-HDBW3849E-AS-IL).
    """
    prefisso_up = prefisso.upper()
    modelli = set()

    for f in MANUALI_DIR.glob("**/*.pdf"):
        nome = f.name
        nome_up = nome.upper()
        if prefisso_up not in nome_up:
            continue
        # Estrai il codice modello: per file tipo DAHUA_IPC-HDBW3849E-AS-IL_Datasheet_...
        # il codice è la seconda parte dopo il primo underscore
        parti = nome.split("_")
        marca = parti[0].upper() if parti else ""
        if marca in ("DAHUA", "RISCO", "ELMO", "NOTIFIER", "ESSER") and len(parti) >= 2:
            codice = parti[1]
        else:
            codice = parti[0]
        # Mantieni solo il codice principale (rimuovi suffissi dopo il 3° trattino)
        trattini = codice.split("-")
        if len(trattini) > 4:
            codice = "-".join(trattini[:4])
        if len(codice) >= 6 and prefisso_up in codice.upper():
            modelli.add(codice)
        if len(modelli) >= max_risultati:
            break

    return sorted(modelli)[:max_risultati]

def trova_e_leggi_pdf(codice_modello: str) -> tuple[str, list]:
    """
    Cerca TUTTI i PDF del modello nella cartella manuali e ne estrae il testo.
    Per prodotti EL.MO. cerca anche i manuali di serie (MT, MU, QG, FT).
    Prioritizza i manuali tecnici MT per domande di programmazione.
    Ritorna (testo_combinato, lista_nomi_file) o ("", []) se non trovato.
    """
    codice_upper = codice_modello.upper()
    # Risolvi alias (es. PROXIMA80 -> PRX80)
    if codice_upper in ALIAS_ELMO:
        codice_upper = ALIAS_ELMO[codice_upper]
        logger.info(f"Alias risolto: {codice_modello.upper()} -> {codice_upper}")

    pdf_trovati = []
    seen = set()

    tutti_file = (list(MANUALI_DIR.glob("**/*.pdf")) +
                  list(MANUALI_DIR.glob("**/*.txt")) +
                  list(MANUALI_DIR.glob("**/*.pptx")) +
                  list(MANUALI_DIR.glob("**/*.PPTX")) +
                  list(MANUALI_DIR.glob("**/*.docx")) +
                  list(MANUALI_DIR.glob("**/*.DOCX")))

    # 1. Cerca corrispondenza diretta nel nome file
    for pdf_path in tutti_file:
        if codice_upper in pdf_path.name.upper() and pdf_path not in seen:
            pdf_trovati.append(pdf_path)
            seen.add(pdf_path)

    # 1b. Se non trovato, prova varianti con prefisso Dahua (es. HDW3849 -> IPC-HDW3849)
    if not pdf_trovati:
        PREFISSI_ESPANSIONE = ["IPC-", "DHI-", "DH-"]
        for prefisso in PREFISSI_ESPANSIONE:
            variante = prefisso + codice_upper
            for pdf_path in tutti_file:
                if variante in pdf_path.name.upper() and pdf_path not in seen:
                    pdf_trovati.append(pdf_path)
                    seen.add(pdf_path)
            if pdf_trovati:
                logger.info(f"[DIAG3] Trovato con variante '{variante}': {len(pdf_trovati)} PDF")
                break

    # 2. Per prodotti EL.MO. cerca anche manuali di famiglia/serie
    # Cerca nel dizionario FAMIGLIE_ELMO le keyword associate al codice
    keywords_famiglia = []
    for chiave, keywords in FAMIGLIE_ELMO.items():
        if codice_upper == chiave or codice_upper.startswith(chiave):
            keywords_famiglia.extend(keywords)
            break
    # Cerca anche per match parziale (es. PRX80 matcha PRX)
    if not keywords_famiglia:
        for chiave, keywords in FAMIGLIE_ELMO.items():
            if chiave in codice_upper or codice_upper in chiave:
                keywords_famiglia.extend(keywords)

    if keywords_famiglia:
        for pdf_path in tutti_file:
            if pdf_path in seen:
                continue
            nome = pdf_path.name.upper()
            for kw in keywords_famiglia:
                if kw.upper() in nome:
                    pdf_trovati.append(pdf_path)
                    seen.add(pdf_path)
                    break

    # 3. Per prodotti RISCO cerca anche manuali di famiglia
    codice_risco = ALIAS_RISCO.get(codice_upper, codice_upper)
    keywords_risco = FAMIGLIE_RISCO.get(codice_risco, [])
    if not keywords_risco:
        # Cerca match parziale nelle famiglie RISCO
        for chiave, keywords in FAMIGLIE_RISCO.items():
            if chiave in codice_upper or codice_upper in chiave:
                keywords_risco.extend(keywords)

    if keywords_risco:
        for pdf_path in tutti_file:
            if pdf_path in seen:
                continue
            nome = pdf_path.name
            if not nome.upper().startswith("RISCO_"):
                continue
            for kw in keywords_risco:
                if kw.lower() in nome.lower():
                    pdf_trovati.append(pdf_path)
                    seen.add(pdf_path)
                    break

    if not pdf_trovati:
        return "", []

    codice_risco = ALIAS_RISCO.get(codice_upper, codice_upper)
    codice_elmo = ALIAS_ELMO.get(codice_upper, codice_upper)
    prodotti_pdf = list({
        codice_upper, codice_risco, codice_elmo,
        risolvi_prodotto_colloquiale(codice_upper),
    })
    domanda_vuota = ""

    prefisso = pdf_trovati[0].name.upper()
    if prefisso.startswith("RISCO_"):
        filtrati = [p for p in pdf_trovati if not _is_risco_pdf_accessorio(p.name)]
        if filtrati:
            pdf_trovati = filtrati
    elif prefisso.startswith("ELMO_"):
        pdf_trovati = _filtra_pdf_elmo(pdf_trovati, prodotti_pdf, domanda_vuota)
    elif pdf_trovati[0].name.upper().startswith("NOTIFIER_"):
        pdf_trovati = _filtra_pdf_notifier(pdf_trovati, prodotti_pdf, domanda_vuota)
    elif pdf_trovati[0].name.upper().startswith("DAHUA_"):
        pdf_trovati = _filtra_pdf_dahua(pdf_trovati, domanda_vuota)

    def priorita(p):
        nome = p.name.upper()
        if nome.startswith("RISCO_"):
            return _priorita_pdf_risco(p, prodotti_pdf)
        if nome.startswith("ELMO_"):
            return _priorita_pdf_elmo(p, prodotti_pdf)
        if nome.startswith("NOTIFIER_"):
            return _priorita_pdf_notifier(p, prodotti_pdf, domanda_vuota)
        if nome.startswith("ESSER_"):
            return _priorita_pdf_esser(p, domanda_vuota)
        if nome.startswith("DAHUA_"):
            return _priorita_pdf_dahua(p, domanda_vuota)
        return (5, p.name.lower())

    pdf_trovati.sort(key=priorita)
    logger.info(f"[DIAG3] trova_e_leggi_pdf('{codice_modello}'): {len(pdf_trovati)} trovati: {[p.name for p in pdf_trovati[:5]]}")

    # Sanity check: se il codice è chiaramente Dahua (HDW/HFW/IPC/NVR ecc.)
    # non includere mai PDF di altre marche nel risultato diretto
    PREFISSI_DAHUA_HARD = ("HDW","HFW","HDB","HFE","HFS","HMW","IPC","NVR","XVR","DHI","HAC","HAD","TPC","PTZ","SD1","SD2","SD3","SD4","SD5","SD6","SD7","SD8","SD9")
    if any(codice_upper.startswith(p) or codice_upper.startswith("IPC-"+p) for p in PREFISSI_DAHUA_HARD):
        pdf_trovati_filtrati = [p for p in pdf_trovati if p.name.upper().startswith("DAHUA_")]
        if pdf_trovati_filtrati:
            logger.info(f"[DIAG3] Filtro hard Dahua: da {len(pdf_trovati)} a {len(pdf_trovati_filtrati)} PDF")
            pdf_trovati = pdf_trovati_filtrati
        elif pdf_trovati:
            logger.warning(f"[DIAG3] Filtro hard Dahua: nessun PDF Dahua trovato, scarto tutti i non-Dahua")
            pdf_trovati = []

    # Leggi fino a 4 PDF (max 15000 caratteri totali)
    testo_totale = ""
    nomi = []
    for pdf_path in pdf_trovati[:4]:
        testo = estrai_testo_file(pdf_path)
        if testo:
            testo_totale += f"\n=== {pdf_path.name} ===\n{testo[:4000]}\n"
            nomi.append(pdf_path.name)
        if len(testo_totale) > 15000:
            break

    # Se non trovato localmente, tenta download online da Dahua
    if not testo_totale:
        logger.info(f"Modello '{codice_modello}' non trovato localmente, tento download online...")
        prod = cerca_prodotto_dahua_online(codice_modello)
        if prod:
            product_id = prod.get("product_id", "")
            prod_nome  = prod.get("product_name", codice_modello)
            for tipo in ["2", "1"]:  # prima manuale, poi datasheet
                testo, nome = scarica_pdf_dahua_online(product_id, tipo, prod_nome)
                if testo:
                    testo_totale += f"\n=== {nome} (scaricato online) ===\n{testo[:6000]}\n"
                    nomi.append(nome)
                    if len(testo_totale) > 12000:
                        break

    return testo_totale, nomi

def cerca_nei_manuali(domanda: str, top_k: int = TOP_K) -> list[dict]:
    """Recupera i chunk più rilevanti per la domanda tramite ricerca semantica."""
    import re
    domanda_up = domanda.upper()
    domanda_norm = re.sub(r'([A-Z]+)\s+([0-9]+[A-Z]*)', r'\1\2', domanda_up)
    codici = list(set(
        re.findall(r'[A-Z0-9]{3,}(?:[-][A-Z0-9]+)+|[A-Z]{3,}[0-9]+[A-Z0-9]*|[A-Z]{2,}[0-9]{2,}[A-Z0-9]*', domanda_up) +
        re.findall(r'[A-Z0-9]{3,}(?:[-][A-Z0-9]+)+|[A-Z]{3,}[0-9]+[A-Z0-9]*|[A-Z]{2,}[0-9]{2,}[A-Z0-9]*', domanda_norm)
    ))
    ESCLUDI2 = {"THE", "AND", "FOR", "CON", "PER", "DAL", "DEL", "UNA", "UNO",
                "NON", "CHE", "SUI", "SUL", "ALL", "SIA", "VIA", "ORA", "GHZ",
                "POI", "BUS", "LED", "LAN", "USB", "SIM", "GSM", "LTE", "POE",
                "TCP", "UDP", "FTP", "SSH", "VPN", "DNS", "NTP", "MAC", "RAM",
                "ROM", "CPU", "GPU", "API", "SDK", "CGI", "URL", "HTTP", "HTTPS",
                "PDF", "ZIP", "JPG", "PNG", "AVI", "MP4", "FPS", "HDR", "SDR",
                "MANUALE", "MANUALI", "CENTRALE", "CENTRALI", "SENSORE", "SENSORI",
                "RIVELATORE", "TASTIERA", "SIRENA", "TELECAMERA", "TELECAMERE",
                "ZONA", "ZONE", "CODICE", "CODICI", "UTENTE", "UTENTI",
                "INSTALLAZIONE", "PROGRAMMAZIONE", "CONFIGURAZIONE", "PROCEDURA",
                "COME", "COSA", "DOVE", "QUANDO", "QUALE", "QUALI", "QUESTO",
                "QUESTA", "QUELLO", "QUELLA", "VOGLIO", "VORREI", "POSSO",
                "PUOI", "RIESCI", "FORNISCI", "DIMMI", "MOSTRAMI", "DAMMI",
                "SERIE", "SISTEMA", "IMPIANTO", "ALLARME", "SEGNALE", "USCITA",
                "INGRESSO", "PARTIZIONE", "AREA", "GRUPPO", "RETE", "CLOUD",
                "RISCO", "DAHUA", "ELMO", "NOTIFIER", "ESSER", "SAMSUNG", "BOSCH", "HONEYWELL",
                "DELLA", "DELLO", "DEGLI", "DELLE", "FORNISCIMI", "CONFIGURO",
                "DIMMI", "MOSTRAMI", "CERCAMI", "TROVAMI", "INSTALLARE",
                "COLLEGARE", "CONNETTERE", "RESETTARE", "PROGRAMMARE",
                "IMPOSTARE", "ABILITARE", "DISABILITARE", "VERIFICARE",
                "CONTROLLARE", "GESTIRE", "AGGIORNARE", "SCARICARE",
                "QUALI", "QUALE", "QUANTO", "QUANTA", "QUANTI", "QUANTE",
                "TUTTO", "TUTTI", "TUTTA", "TUTTE", "PRIMO", "PRIMA",
                "SECONDO", "SECONDA", "NUOVO", "NUOVA", "VECCHIO",
                "GRANDE", "PICCOLO", "ESTERNO", "INTERNO", "WIRELESS",
                "CABLATO", "DIGITALE", "ANALOGICO", "INTEGRATO",
                "VERSIONE", "MODELLO", "PRODOTTO", "DISPOSITIVO",
                "FORNICA", "FORNICIMI", "FORNIRE", "INVIARMI",
                "MANDAMI", "INVIAMI", "CERCAMI", "AIUTAMI", "SPIEGAMI",
                "CABLATO", "FUNZIONE", "FUNZIONA", "FUNZIONAMENTO", 
                "INCENDIO", "FUOCO", "EMERGENZA", "RACCOLTA", "EVACUAZIONE",
                "ESTERNO", "ESTERNO", "INTERNO", "ACCESSO", "VARCO", "INGRESSO",
                "USCITA", "PARCHEGGIO", "CORTILE", "AREA", "PERIMETRO",
                "COMANDI", "COMANDO", "INVIARE", "POSSIBILITA", "HANNO",
                "POSSONO", "POSSIBILE", "INVIO", "TRAMITE", "MEDIANTE",
                "NOTIFICA", "NOTIFICHE", "EVENTO", "EVENTI", "AZIONE", "AZIONI",
                "ALLARMI", "REGOLA", "REGOLE", "TRIGGER", "PORTA", "PORTE",
                "INDIRIZZO", "INDIRIZZI", "RISPOSTA", "RICHIESTA", "RICHIESTE",
                "CHIAMATA", "CHIAMATE", "METODO", "METODI", "PROTOCOLLO",
                "INTEGRAZIONE", "FUNZIONALITA", "SUPPORTO", "SUPPORTA"}
    codici = [c for c in codici if c not in ESCLUDI2 and c not in PAROLE_NON_CODICI and codice_modello_plausibile(c)]

    # Filtro post-query: quando la domanda menziona una famiglia specifica,
    # dopo la ricerca semantica teniamo solo i chunk dei file giusti.
    # NON usiamo where=$contains perché ChromaDB 1.5.7 non lo supporta su stringhe.
    domanda_up_search = domanda.upper()
    FILTRI_SPECIFICI = {
        ("PROXIMA", "PRX80", "PRX128", "PRX256", "PRX1024"): "MT_serie-PROXIMA_2025",
        ("SUPERIA", "SPR256", "SPR512"):                      "SUPERIA-series",
        ("PREGIO", "PREGIO500", "PREGIO2000"):                "PREGIO",
        ("VIDOMO", "VIDOMO2K"):                               "VIDOMO",
        ("KARMA",):                                           "KARMA",
        # Dahua — staffaggi / accessori (evita chunk ESSER/altre marche in ricerca semantica)
        ("HFW", "HDW", "HDB", "HAC", "HAD", "IPC-H", "PFA150", "PFB2203", "PFA130",
         "CAMERA-ACCESSORIES", "ACCESSORIES-SELECTION"):      "Camera-Accessories-Selection",
        # RISCO — famiglie (ordine: varianti specifiche prima di LIGHTSYS generico)
        ("LIGHTSYS2", "LIGHT SYS 2", "LIGHTSYS 2"):           "LightSYS-2-Full-Installation-Manual",
        ("LIGHTSYS+", "LIGHTSYSPLUS", "LIGHT SYS PLUS"):       "LightSYS-Plus-Installer-Manual-IT",
        ("LIGHTSYSAIR", "LIGHT SYS AIR"):                     "LightSYS_Air",
        ("AGILITY4", "AGILITY 4"):                            "Agility_4",
        ("AGILITY",):                                         "Agility",
        ("WICOMM", "WI COMM"):                                "WiComm",
        ("LIGHTSYS", "LIGHT SYS"):                            "lightsys-2",
        # NOTIFIER
        ("AM-8200N", "AM8200N", "AM-8200"):                   "AM-8200N manu-prog",
        ("AM-8100", "AM8100"):                                "AM-8100 manu-prog",
        ("NFS2", "NFS-2"):                                    "NFS2",
        ("ID3000",):                                          "ID3000",
        ("M5A", "W5A", "M3A", "CALL POINT"):                  "M5A-W5A",
        ("FSP", "FDM",):                                      "FSP",
        # ESSER
        ("FLEXES", "FLEX ES", "FLEX CONTROL"):                "FlexES",
        ("IQ8", "IQ8CONTROL", "ESSERBUS", "8000C"):           "esserbus",
        ("WINMAG",):                                          "Winmag",
        # ELMO — estensione
        ("HEKLA", "HLNET"):                                   "HEKLA",
        ("HERCOLA",):                                         "HERCOLA",
        ("TACORA", "TA2000"):                                 "TACORA",
        ("GATEWAY2K", "GATEWAY 2K", "NG-TRX", "NG TRX"):      "MT_GATEWAY2K",
        ("ECONNECT", "E-CONNECT"):                            "e-Connect",
    }
    file_filter = None
    for keywords_tuple, pattern in FILTRI_SPECIFICI.items():
        if any(kw in domanda_up_search for kw in keywords_tuple):
            file_filter = pattern
            break

    risultati = collection.query(
        query_texts=[domanda],
        n_results=min(top_k * 3 if file_filter else top_k, collection.count()),
        include=["documents", "metadatas", "distances"],
    )
    chunks = []
    if risultati and risultati["documents"]:
        for doc, meta, dist in zip(
            risultati["documents"][0],
            risultati["metadatas"][0],
            risultati["distances"][0],
        ):
            score = round(1 - dist, 3)
            fonte = meta.get("source", "?")
            for codice in codici:
                if codice in fonte.upper():
                    score = min(1.0, score + 0.3)
                    break
            chunks.append({
                "testo":  doc,
                "fonte":  fonte,
                "pagina": meta.get("pagina", "?"),
                "score":  score,
            })

    # Post-filtro per famiglia: se ci sono chunk dal file giusto, usa solo quelli
    if file_filter:
        chunks_filtrati = [c for c in chunks if file_filter.upper() in c["fonte"].upper()]
        if chunks_filtrati:
            logger.info(f"Post-filtro '{file_filter}': {len(chunks_filtrati)}/{len(chunks)} chunk mantenuti")
            chunks = chunks_filtrati[:top_k]
        else:
            logger.info(f"Post-filtro '{file_filter}': nessun match, uso tutti i chunk")
            chunks = chunks[:top_k]
    else:
        chunks = chunks[:top_k]

    chunks.sort(key=lambda x: x["score"], reverse=True)
    return chunks

def pulisci_testo(testo: str) -> str:
    """Rimuove caratteri di formattazione non compatibili con TTS."""
    for char in ["*", "#", "_", "~", "`", ">", "[", "]", "(", ")", "•", "–", "—"]:
        testo = testo.replace(char, "")
    return testo

def is_richiesta_pdf(domanda: str) -> bool:
    """Rileva se l'utente chiede esplicitamente il file PDF o il manuale."""
    keywords = [
        # Richieste esplicite file
        "pdf", "file", "documento", "allegato", "scarica", "invia",
        "mandami", "puoi inviarmi", "puoi mandarmi", "fornisci",
        "forniscimi", "download", "datasheet", "scheda tecnica",
        "per mail", "via mail", "via email", "allegare",
        # Richieste manuale (tipiche anche da vocale)
        "mandami il manuale", "inviami il manuale", "forniscimi il manuale",
        "puoi mandarmi il manuale", "voglio il manuale", "dammi il manuale",
        "mandami la scheda", "inviami la scheda", "voglio la scheda",
        "mandami il datasheet", "inviami il datasheet",
    ]
    d = domanda.lower()
    return any(k in d for k in keywords)

def is_domanda_dimensionamento(domanda: str, codici_trovati: list = None,
                              marche_attese: set = None) -> bool:
    """Opzione B: carica le linee guida per domande tecniche generiche Dahua/dimensionamento."""
    domanda_low = domanda.lower().strip()
    NON_TECNICI = {"ciao", "salve", "buongiorno", "buonasera", "buon pomeriggio",
                   "grazie", "ok", "perfetto", "capito", "ricevuto", "ottimo",
                   "bene", "va bene", "ho capito", "inteso"}
    if domanda_low in NON_TECNICI or len(domanda_low) < 8:
        return False
    # Non caricare linee guida Dahua per domande su altre marche
    if marche_attese and marche_attese - {"DAHUA"}:
        return False
    MARCA_NON_DAHUA = [
        "risco", "lightsys", "agility", "wicomm", "elmo", "el.mo", "notifier",
        "esser", "proxima", "iq8", "flexes", "passlight", "superia", "pregio",
    ]
    if any(m in domanda_low for m in MARCA_NON_DAHUA):
        return False
    # Domande operative (antintrusione, antincendio, video) non sono dimensionamento Dahua
    if any(k in domanda_low for k in KW_OPERATIVE_TUTTE):
        return False
    if codici_trovati:
        SELEZIONE = ["va bene", "adatto", "posso usare", "posso installare",
                     "quale", "che ", "consiglia", "suggerisci", "differenza",
                     "confronto", "alternativa", "invece di", "al posto",
                     "meglio", "oppure", "compatibile"]
        if not any(k in domanda_low for k in SELEZIONE):
            return False
    return True

def _is_domanda_dimensionamento_UNUSED(domanda: str) -> bool:
    """Versione originale conservata per riferimento - non usata."""
    keywords = [
        "impianto", "dimensionamento", "progetto", "installaz",
        "telecamere da installare", "quante telecamere", "sistema",
        "proposta", "cosa consigli", "cosa posso proporre",
        "configurazione", "preventivo", "devo fare", "devo installare",
        "conteggio persone", "people counting", "contare persone",
        "quante persone", "zona di raccolta", "punto di raccolta",
        "monitorare una zona", "telecamere per", "modelli di telecamere",
        "quali modelli", "quale telecamera", "quale modello",
        "mi serve una telecamera", "mi servono telecamere",
        "telecamera per esterni", "telecamera esterna",
        # Staffaggi e accessori
        "staffa", "staffe", "accessori", "montaggio", "bracket",
        "junction box", "come si monta", "come installo",
        "come si installa", "supporto", "fissaggio",
        # ANPR / targhe / varco
        "targa", "targhe", "anpr", "varco", "carraio", "cancello",
        "riconoscimento targa", "lettura targa", "white list", "blacklist",
        "black list", "lista targhe", "accesso veicoli", "veicolo",
        "barriera", "sbarra", "parcheggio",
        # Termiche / speciali
        "termica", "termiche", "temperatura", "calore",
        # PTZ / deterrenza
        "deterrenza", "tioc", "full color", "colorvu",
        "sirena integrata", "luce integrata",
        # Domande tecniche avanzate che richiedono il mapping manuali
        "http", "cgi", "api", "sdk", "webhook", "integrazione",
        "comando", "comandi", "notifica http", "action url",
        "event", "trigger url", "alarm action",
        # Notifier / antincendio
        "pulsante", "rottura vetro", "call point", "rivelatore", "centrale",
        "incendio", "fumo", "calore", "termico", "antincendio",
        "cella frigo", "cella frigorifera", "ip67", "ip65",
        "notifier", "nfs", "m3a", "m5a", "w5a", "bg6", "id3000",
        "indirizzato", "convenzionale", "quale pulsante", "che pulsante",
        "bassa temperatura", "ambiente umido", "ambiente esterno",
    ]
    return any(k in domanda.lower() for k in keywords)

def carica_linee_guida() -> str:
    """Carica il file delle linee guida se presente."""
    percorsi = [
        MANUALI_DIR / "linee_guida_dimensionamento.txt",
        Path("./manuali/linee_guida_dimensionamento.txt"),
        Path("./linee_guida_dimensionamento.txt"),
    ]
    for p in percorsi:
        if p.exists():
            try:
                return p.read_text(encoding="utf-8")
            except Exception:
                pass
    return ""


CODICE_MODELLO_DAHUA_RE = re.compile(
    r"\b(?:IPC|HAC|DH|DHI|NVR|XVR|SD|TPC|PTZ)-[\w][\w\-]*\b",
    re.IGNORECASE,
)


def _normalizza_codice_dahua(codice: str) -> str:
    c = codice.upper().strip().replace(" ", "")
    for pref in ("IPC-", "DH-", "DHI-", "HAC-"):
        if c.startswith(pref):
            c = c[len(pref):]
    return c


def _varianti_codice_flat(codice: str) -> list[str]:
    """Varianti tolleranti (es. ASPV nelle linee guida = AS-PV nel datasheet)."""
    base = _normalizza_codice_dahua(codice).replace("-", "").upper()
    if len(base) < 6:
        return []
    varianti = {base}
    sostituzioni = (
        ("ASPV", "ASPVPV"),
        ("ZASPV", "ZASPVPV"),
        ("ASPVPRO", "ASPVPVPRO"),
        ("ZASPVPRO", "ZASPVPVPRO"),
    )
    for src, dst in sostituzioni:
        if src in base:
            varianti.add(base.replace(src, dst, 1))
        if dst in base:
            varianti.add(base.replace(dst, src, 1))
    return list(varianti)


def _codice_compatibile_filename(codice: str, nome_file: str) -> bool:
    """Match tollerante (es. ASPV vs AS-PV) tra codice e nome PDF."""
    nf = nome_file.upper().replace("-", "").replace("_", "")
    for cc in _varianti_codice_flat(codice):
        if cc in nf:
            return True
    return False


def estrai_codici_modello_da_testo(testo: str) -> list[str]:
    trovati = []
    visti = set()
    for m in CODICE_MODELLO_DAHUA_RE.finditer(testo):
        cod = m.group(0).upper().strip()
        if len(cod) < 10 or cod in visti:
            continue
        visti.add(cod)
        trovati.append(cod)
    return trovati


def seleziona_sezioni_linee_guida(domanda: str, testo: str, max_chars: int = 5500) -> str:
    """Estrae dalle linee guida solo le sezioni pertinenti alla domanda."""
    if not testo:
        return ""
    domanda_low = domanda.lower()
    parole_domanda = set(re.findall(r"[a-z0-9]{4,}", domanda_low))

    keywords_extra = set()
    if any(k in domanda_low for k in ("eyeball", "dome", "cupola", "minidome")):
        keywords_extra.update({"eyeball", "dome", "hdw", "tioc", "cupola"})
    if any(k in domanda_low for k in ("bullet", "tubolare", "hfw")):
        keywords_extra.update({"bullet", "hfw", "tubolare"})
    if any(k in domanda_low for k in ("audio", "bidirez", "microfono", "altoparlante", "speaker", "parlare")):
        keywords_extra.update({"audio", "deterrenza", "tioc", "as-pv", "aspv", "tm-ase", "speaker", "microfono"})
    if any(k in domanda_low for k in ("staffa", "montaggio", "pfb", "pfa")):
        keywords_extra.update({"staffa", "staffaggio", "pfb", "pfa", "montaggio"})
    if any(k in domanda_low for k in ("anpr", "targa", "varco")):
        keywords_extra.update({"anpr", "targa", "varco"})
    parole_domanda |= keywords_extra

    blocchi = re.split(r"(?=^={3,}|^--- )", testo, flags=re.MULTILINE)
    if len(blocchi) <= 1:
        return testo[:max_chars]

    header = blocchi[0].strip()
    sezioni = []
    for blocco in blocchi[1:]:
        blocco_low = blocco.lower()
        score = sum(1 for p in parole_domanda if p in blocco_low)
        codici = estrai_codici_modello_da_testo(blocco)
        if score > 0 or codici:
            sezioni.append((score, blocco.strip()))

    sezioni.sort(key=lambda x: x[0], reverse=True)
    if not sezioni:
        return testo[:max_chars]

    out = header + "\n\n" if header else ""
    for score, blocco in sezioni:
        if score <= 0 and len(out) > 500:
            continue
        if len(out) + len(blocco) + 2 > max_chars:
            break
        out += blocco + "\n\n"
    return out.strip() or testo[:max_chars]


def _modello_dahua_pertinente_domanda(codice: str, domanda_low: str) -> bool:
    c = codice.upper()
    if any(k in domanda_low for k in ("eyeball", "dome", "cupola")) and not any(
        k in domanda_low for k in ("bullet", "tubolare")
    ):
        if "HFW" in c or (c.startswith("IPC-HDB") and "HDBW" not in c):
            return False
        if "HDW" not in c and "HDBW" not in c:
            return False
    if any(k in domanda_low for k in ("bullet", "tubolare")) and "HFW" not in c:
        return False
    return True


def _priorita_modello_dimensionamento(codice: str, domanda_low: str) -> tuple:
    c = codice.upper()
    score = 0
    if any(k in domanda_low for k in ("audio", "bidirez", "microfono", "altoparlante", "speaker")):
        if "AS-PV" in c or "ASPV" in c or "ZAS-PV" in c or "ZASPV" in c:
            score += 5
        if "TM-ASE" in c or "TM-AS" in c or "-ASE" in c:
            score += 4
        if "TQ" in c and "HDW" in c:
            score += 3
    if any(k in domanda_low for k in ("eyeball", "dome", "cupola")) and "HDW" in c:
        score += 2
    if "TIOC" in domanda_low or "deterrenza" in domanda_low:
        if "AS-PV" in c or "ASPV" in c:
            score += 2
    return (-score, c)


def _codici_datasheet_da_cercare(codice: str) -> list[str]:
    """Codice e varianti senza revisione firmware (es. -S5 nel catalogo, _S5 nel PDF)."""
    cod = codice.upper().strip()
    varianti = [cod]
    base = re.sub(r"-S\d+$", "", cod)
    if base != cod:
        varianti.append(base)
    return varianti


def trova_datasheet_dahua(codice: str) -> Path | None:
    """Cerca il datasheet PDF Dahua per un codice modello."""
    candidati: list[Path] = []
    visti: set[str] = set()

    def _aggiungi(matches: list[Path], cod: str) -> None:
        for pdf in matches:
            if pdf.suffix.lower() != ".pdf":
                continue
            key = str(pdf.resolve())
            if key in visti:
                continue
            nome_up = pdf.name.upper()
            if "DATASHEET" not in nome_up:
                continue
            if _codice_compatibile_filename(cod, pdf.name):
                visti.add(key)
                candidati.append(pdf)

    for cod in _codici_datasheet_da_cercare(codice):
        cc_norm = _normalizza_codice_dahua(cod)
        stem = cod.replace("-", "")
        for pattern in (
            f"DAHUA_{cod}*Datasheet*.pdf",
            f"DAHUA_*{cod}*Datasheet*.pdf",
            f"DAHUA_*{cc_norm}*Datasheet*.pdf",
            f"*{stem}*datasheet*.pdf",
            f"**/*{cod}*Datasheet*.pdf",
        ):
            _aggiungi(list(MANUALI_DIR.glob(pattern)), cod)

    if not candidati:
        return None
    candidati.sort(key=lambda p: (
        0 if p.name.upper().startswith("DAHUA_") else 1,
        -p.stat().st_mtime,
    ))
    return candidati[0]


def estrai_sezioni_rilevanti_datasheet(testo: str, domanda: str, max_chars: int = 2200) -> str:
    """Estrae dal datasheet le parti utili per la domanda (audio, alimentazione, ecc.)."""
    if not testo:
        return ""
    domanda_low = domanda.lower()
    focus = []
    if any(k in domanda_low for k in ("audio", "bidirez", "microfono", "altoparlante", "speaker", "parlare")):
        focus = [
            "audio", "microphone", "mic ", "speaker", "two-way", "two way",
            "intercom", "sound", "bidirez", "talk", "voice",
        ]
    if any(k in domanda_low for k in ("ip67", "ip66", "ip rating", "stagno", "esterno")):
        focus.extend(["ip67", "ip66", "ingress protection", "waterproof", "outdoor"])
    if any(k in domanda_low for k in ("wdr", "starlight", "notturn", "illuminator", "ir ")):
        focus.extend(["wdr", "starlight", "illuminator", "ir ", "lux", "0.00"])

    if not focus:
        return testo[:max_chars]

    righe_utili = []
    for riga in testo.splitlines():
        rl = riga.lower()
        if any(k in rl for k in focus):
            righe_utili.append(riga.strip())
    if not righe_utili:
        return testo[:max_chars]

    estratto = "\n".join(righe_utili)
    if len(estratto) > max_chars:
        return estratto[:max_chars]
    # Aggiungi intestazione modello se troppo corto
    if len(estratto) < 400:
        head = testo[: min(800, max_chars // 2)]
        return (head + "\n...\n" + estratto)[:max_chars]
    return estratto


def carica_datasheet_per_dimensionamento(
    domanda: str,
    linee_guida: str,
    max_modelli: int = 10,
    max_chars_total: int = 12000,
) -> tuple[str, set]:
    """
    Dalle linee guida ricava i candidati e carica i datasheet corrispondenti.
    """
    sezioni = seleziona_sezioni_linee_guida(domanda, linee_guida)
    codici = estrai_codici_modello_da_testo(sezioni)
    if len(codici) < 3:
        codici = list(dict.fromkeys(codici + estrai_codici_modello_da_testo(linee_guida[:20000])))

    domanda_low = domanda.lower()
    codici = [c for c in codici if _modello_dahua_pertinente_domanda(c, domanda_low)]
    codici.sort(key=lambda c: _priorita_modello_dimensionamento(c, domanda_low))

    contesto = ""
    fonti: set = set()
    chars = 0
    per_modello = max(1500, max_chars_total // max(max_modelli, 1))

    for codice in codici:
        if len(fonti) >= max_modelli or chars >= max_chars_total:
            break
        pdf = trova_datasheet_dahua(codice)
        if not pdf:
            logger.info(f"Datasheet non trovato per candidato linee guida: {codice}")
            continue
        testo = estrai_testo_file(pdf)
        if not testo:
            continue
        estratto = estrai_sezioni_rilevanti_datasheet(testo, domanda, per_modello)
        contesto += f"\n=== DATASHEET {codice} ({pdf.name}) ===\n{estratto}\n"
        fonti.add(pdf.name)
        chars += len(estratto)
        logger.info(f"Datasheet dimensionamento caricato: {codice} -> {pdf.name}")

    return contesto, fonti


def cerca_prodotto_dahua_online(codice: str) -> dict:
    """
    Cerca un prodotto sul portale Dahua per codice modello.
    Ritorna dict con product_id, product_name o {} se non trovato.
    """
    try:
        r = req_lib.post(
            DAHUA_API_SEARCH,
            data={"id": "5", "keyword": codice, "page": "1"},
            headers=DAHUA_HEADERS, timeout=15
        )
        data = r.json().get("data", {})
        if isinstance(data, dict):
            prodotti = data.get("product", [])
            # Cerca corrispondenza esatta o parziale
            for p in prodotti:
                nome = p.get("product_name", "").upper()
                if codice.upper() in nome:
                    return p
            if prodotti:
                return prodotti[0]
    except Exception as e:
        logger.debug(f"Errore ricerca Dahua online {codice}: {e}")
    return {}

def scarica_pdf_dahua_online(product_id: str, tipo: str, prod_nome: str) -> tuple[str, str]:
    """
    Scarica un PDF da Dahua online dato il product_id.
    tipo: "1"=Datasheet, "2"=Manuale Utente
    Salva nella cartella manuali e ritorna (testo_estratto, nome_file).
    """
    try:
        # Ottieni URL del file
        r = req_lib.post(
            DAHUA_API_FILZIP,
            data={"id": product_id, "type": tipo},
            headers=DAHUA_HEADERS, timeout=30
        )
        data = r.json()
        if data.get("status") != "00" or not data.get("data"):
            return "", ""

        url = data["data"]
        if "downloading?url=" in url:
            url = url.split("downloading?url=")[1]

        # Scarica il file
        r2 = req_lib.get(url, timeout=60, headers=DAHUA_HEADERS)
        if r2.status_code != 200:
            return "", ""

        ct = r2.headers.get("content-type", "")
        tipo_label = "Datasheet" if tipo == "1" else "Manuale_Utente"
        prod_safe = re.sub(r'[<>:"/\\|?*\s]', "_", prod_nome)[:80]

        # Gestisci ZIP o PDF diretto
        if "zip" in ct or url.endswith(".zip") or "previous.dahuasecurity.com" in url:
            try:
                with zipfile.ZipFile(io.BytesIO(r2.content)) as zf:
                    pdf_files = [n for n in zf.namelist() if n.lower().endswith(".pdf")]
                    if not pdf_files:
                        return "", ""
                    # Prendi il primo PDF
                    pdf_name = pdf_files[0]
                    pdf_bytes = zf.read(pdf_name)
                    base = os.path.basename(pdf_name)
                    nome_file = f"DAHUA_{prod_safe}_{tipo_label}_{base}"
            except zipfile.BadZipFile:
                return "", ""
        elif "pdf" in ct or url.endswith(".pdf"):
            pdf_bytes = r2.content
            nome_url = url.split("/")[-1].split("?")[0]
            nome_file = f"DAHUA_{prod_safe}_{tipo_label}_{nome_url}"
        else:
            return "", ""

        # Salva il file
        dest = MANUALI_DIR / nome_file
        MANUALI_DIR.mkdir(parents=True, exist_ok=True)
        with open(dest, "wb") as f:
            f.write(pdf_bytes)
        logger.info(f"✅ Scaricato online e salvato: {nome_file}")

        # Estrai testo
        testo = estrai_testo_pdf(dest)

        # Indicizza in background
        import threading
        def indicizza_nuovo():
            try:
                file_hash = hash_file(dest)
                chunks = chunking(testo)
                doc_id_prefix = f"{dest.stem}_{file_hash}"
                ids = [f"{doc_id_prefix}_chunk{i}" for i in range(len(chunks))]
                metadatas = [{"source": nome_file, "file_hash": file_hash,
                              "chunk_idx": i, "pagina": f"~{i+1}"}
                             for i in range(len(chunks))]
                for b in range(0, len(chunks), 500):
                    collection.upsert(ids=ids[b:b+500], documents=chunks[b:b+500],
                                      metadatas=metadatas[b:b+500])
                logger.info(f"✅ Indicizzato: {nome_file}")
            except Exception as e:
                logger.error(f"Errore indicizzazione {nome_file}: {e}")
        threading.Thread(target=indicizza_nuovo, daemon=True).start()

        return testo, nome_file

    except Exception as e:
        logger.error(f"Errore download Dahua online {product_id}: {e}")
        return "", ""

def risolvi_prodotto_colloquiale(prodotto: str) -> str:
    """Applica alias marca (es. PROXIMA → PRX, LIGHTSYSPLUS → LIGHTSYS+)."""
    p = prodotto.upper().strip()
    return ALIAS_RISCO.get(p, ALIAS_ELMO.get(p, p))


def estrai_prodotti_da_testo(testo: str) -> list:
    """
    Estrae nomi colloquiali di prodotti/famiglie da domande in linguaggio naturale
    (tutte le marche in archivio: RISCO, ELMO, NOTIFIER, ESSER, DAHUA).
    """
    if not testo or not testo.strip():
        return []
    low = testo.lower()
    up = testo.upper()
    up = up.replace("LIGHTSYS+", "LIGHTSYSPLUS")
    up = up.replace("LIGHTSYS PLUS", "LIGHTSYSPLUS")
    up = up.replace("LIGHTSYS AIR", "LIGHTSYSAIR")
    up = up.replace("LIGHTSYS 2", "LIGHTSYS2")
    up = up.replace("WICOMM PRO", "WICOMMPRO")
    up = up.replace("AGILITY 4", "AGILITY4")
    up = up.replace("GATEWAY 2K", "GATEWAY2K")

    trovati = []
    for needle, chiave in PATTERN_PRODOTTI_COLLOQUIALI:
        if needle in low and chiave not in trovati:
            trovati.append(chiave)

    for chiave in sorted(FAMIGLIE_RISCO.keys(), key=len, reverse=True):
        if chiave in up and chiave not in trovati:
            trovati.append(chiave)
    for chiave in sorted(FAMIGLIE_ELMO.keys(), key=len, reverse=True):
        if chiave in up and chiave not in trovati:
            trovati.append(chiave)

    return _normalizza_prodotti_estratti(trovati)


def _normalizza_prodotti_estratti(trovati: list) -> list:
    """Rimuove prodotti generici quando è già presente una variante più specifica."""
    if not trovati:
        return trovati
    t = list(trovati)
    # RISCO
    if any(p in t for p in ("LIGHTSYS+", "LIGHTSYSPLUS", "LIGHTSYS2", "LIGHTSYSAIR")):
        t = [x for x in t if x != "LIGHTSYS"]
    if "LIGHTSYS+" in t and "LIGHTSYSPLUS" in t:
        t = [x for x in t if x != "LIGHTSYSPLUS"]
    # ELMO
    if "PREGIO2000" in t:
        t = [x for x in t if x not in ("PREGIO", "PREGIO500")]
    elif "PREGIO500" in t:
        t = [x for x in t if x != "PREGIO"]
    for prx in ("PRX1024", "PRX256", "PRX128", "PRX80"):
        if prx in t:
            t = [x for x in t if x not in ("PRX", "PROXIMA")]
    if any(x.startswith("PRX") for x in t) and "PROXIMA" in t:
        t = [x for x in t if x != "PROXIMA"]
    if "VIDOMO2K" in t:
        t = [x for x in t if x != "VIDOMO"]
    if "GATEWAY2K" in t:
        t = [x for x in t if x != "GATEWAY"]
    if "SPR2040" in t:
        t = [x for x in t if x not in ("SPR", "SUPERIA")]
    elif any(x in t for x in ("SPR512", "SPR256")) and "SUPERIA" in t:
        t = [x for x in t if x != "SUPERIA"]
    return t


def _menziona_marca(domanda_low: str, marca: str) -> bool:
    return any(k in domanda_low for k in MARCA_MENTION_KEYWORDS.get(marca, []))


def _domanda_operativa_per_marca(domanda_low: str, marca: str) -> bool:
    if marca in ("RISCO", "ELMO"):
        return any(k in domanda_low for k in KW_OPERATIVE_ANTINTRUSIONE)
    if marca in ("NOTIFIER", "ESSER"):
        return any(k in domanda_low for k in KW_OPERATIVE_ANTINCENDIO)
    if marca == "DAHUA":
        return any(k in domanda_low for k in KW_OPERATIVE_VIDEO)
    return False


def _marche_da_domanda_colloquiale(domanda: str, prodotti: list, marche_attese: set | None) -> set:
    """Determina quali marche processare per caricamento manuale mirato."""
    domanda_low = domanda.lower()
    if marche_attese:
        return set(marche_attese)

    marche = {m for m in MARCA_MENTION_KEYWORDS if _menziona_marca(domanda_low, m)}
    if marche:
        return marche

    def _senza_altra_marca(escluse_marche: set) -> bool:
        for m, kws in MARCA_MENTION_KEYWORDS.items():
            if m in escluse_marche:
                continue
            if any(k in domanda_low for k in kws):
                return False
        return True

    if any(k in domanda_low for k in KW_OPERATIVE_ANTINTRUSIONE):
        if _senza_altra_marca({"RISCO", "ELMO"}):
            marche.add("RISCO")
    if any(k in domanda_low for k in KW_OPERATIVE_ANTINCENDIO):
        if _menziona_marca(domanda_low, "ESSER"):
            marche.add("ESSER")
        elif _senza_altra_marca({"NOTIFIER", "ESSER"}):
            marche.add("NOTIFIER")
    if any(k in domanda_low for k in KW_OPERATIVE_VIDEO):
        if _senza_altra_marca({"DAHUA"}):
            marche.add("DAHUA")
    return marche


# PDF RISCO periferici/accessori (non manuali centrale/programmazione)
RISCO_PDF_ACCESSORIO_MARKERS = (
    "power-supply", "power_supply", "plastic_3.5a", "expander",
    "communication-module", "4g-module", "box_4g", "brochure",
    "ip-module", "keypad-instructions", "panda-for", "axesplus-job",
)


def _domanda_risco_accessorio(domanda_low: str) -> bool:
    """True se la domanda riguarda esplicitamente un accessorio, non la centrale."""
    return any(
        k in domanda_low
        for k in (
            "alimentatore", "power supply", "espansione radio", "modulo 4g",
            "modulo ip", "tastiera panda", "brochure", "rp432eps",
        )
    )


def _is_risco_pdf_accessorio(nome_file: str) -> bool:
    n = nome_file.lower()
    return any(m in n for m in RISCO_PDF_ACCESSORIO_MARKERS)


def _prodotti_risco_plus(prodotti: list | None) -> bool:
    if not prodotti:
        return False
    up = {p.upper() for p in prodotti}
    return bool(up & {"LIGHTSYS+", "LIGHTSYSPLUS"})


def _priorita_pdf_risco(path: Path, prodotti: list | None = None) -> tuple:
    """Ordina manuali RISCO: centrale/programmazione IT prima di accessori."""
    n = path.name.lower()
    if _is_risco_pdf_accessorio(path.name):
        return (9, n)

    plus = _prodotti_risco_plus(prodotti) or "lightsys_plus" in n or "lightsys-plus" in n or "lightsys+" in n
    if plus:
        if "installer-manual" in n and "-it-" in n:
            return (0, n)
        if "full-user" in n and "-it-" in n:
            return (1, n)
        if "quick-installer-guide" in n and "-it-" in n:
            return (2, n)
        if "installer-manual" in n:
            return (3, n)
        if "full-user" in n or "user-manual" in n:
            return (4, n)
        if "quick-user" in n:
            return (5, n)
        return (6, n)

    if "lightsys-2" in n or "lightsys2" in n:
        if "full-installation-manual" in n and "-it-" in n:
            return (0, n)
        if "full-installation-manual" in n:
            return (1, n)
        if "quick-installation" in n:
            return (4, n)
        return (3, n)

    if "full-installation-manual" in n and "-it-" in n:
        return (0, n)
    if "full-installation-manual" in n:
        return (1, n)
    if "installer-manual" in n and "-it-" in n:
        return (2, n)
    if "installation" in n or "program" in n:
        return (3, n)
    return (5, n)


def _cerca_pdf_risco_pattern(pattern: str, prodotti: list, domanda_low: str) -> list:
    """Cerca PDF RISCO provando pattern principali e alias LightSYS Plus."""
    patterns = [pattern]
    if "LightSYS" in pattern and "Plus" in pattern:
        patterns = [
            "LightSYS-Plus-Installer-Manual-IT",
            "LightSYS-Plus-Full-User-IT",
            "LightSYS_Plus_Installer",
            "LightSYS_Plus_Full_User",
            "LightSYS_Plus",
            "lightsys-Plus",
        ]
    visti = set()
    matches = []
    for pat in patterns:
        for p in cerca_file_per_mappa(pat, {"RISCO"}, (".pdf",)):
            if p not in visti:
                visti.add(p)
                matches.append(p)
    if not _domanda_risco_accessorio(domanda_low):
        centrale = [p for p in matches if not _is_risco_pdf_accessorio(p.name)]
        if centrale:
            matches = centrale
    return matches


# ── ELMO: fogli tecnici / schede periferiche vs manuali serie (MT/MP/MU) ──
ELMO_PDF_ACCESSORIO_MARKERS = (
    "_ft_", "_dc_", "_ft.", "_dc.", "_qg_", "sctlt", "modbus_proxima",
    "_ce_", "_imq",  # certificazioni CE/IMQ, non manuali tecnici
)


def _domanda_elmo_accessorio(domanda_low: str) -> bool:
    return any(k in domanda_low for k in ("modbus", "foglio tecnico", "datasheet elmo", "scheda tecnica"))


def _is_elmo_pdf_accessorio(nome_file: str, domanda_low: str = "") -> bool:
    n = nome_file.lower()
    # MT/MP/MU restano manuali tecnici anche se nel nome compare IMQ/CE (certificazione)
    if "_mt_" in n or n.startswith("mt_") or "_mp_" in n or "_mu_" in n:
        return False
    if _domanda_elmo_accessorio(domanda_low) and ("_ft_" in n or "_dc_" in n):
        return False
    if any(m in n for m in ("_ft_", "_dc_", "_ft.", "_dc.", "_qg_", "sctlt", "modbus_proxima")):
        return True
    if "_doc" in n or "libretto" in n:
        return True
    if "_ce_" in n or "_imq" in n:
        return True
    return False


def _pattern_aliases_elmo(pattern: str, prodotti: list) -> list:
    aliases = [pattern]
    if "PROXIMA" in pattern or "serie-PROXIMA" in pattern:
        aliases = [
            "MT_serie-PROXIMA_2025-06", "MT_serie-PROXIMA", "serie-PROXIMA_MT",
            "serie-PROXIMA", "MP_serie-PROXIMA",
        ]
    elif pattern.upper() in ("PREGIO500", "PREGIO2000"):
        sk = "SKPREGIO500" if "500" in pattern.upper() else "SKPREGIO2000"
        aliases = [pattern, sk, "MT_serie-PREGIO", "serie-PREGIO", "Pregio-series"]
    elif "PREGIO" in pattern:
        aliases = ["MT_serie-PREGIO", "serie-PREGIO", "Pregio-series", "PREGIO"]
    elif "SUPERIA" in pattern:
        aliases = ["SUPERIA-series", "TM_SUPERIA", "SUPERIA_MT", "SUPERIA"]
    elif "VIDOMO" in pattern:
        aliases = ["VIDOMO2K", "VILLEGGIO-NG", "VIDOMO", "VILLEGGIO"]
    elif "GATEWAY" in pattern:
        aliases = ["MT_GATEWAY2K", "GATEWAY2K_MT", "GATEWAY2K", "MT_GATEWAY"]
    elif "HEKLA" in pattern or "HLNET" in pattern:
        aliases = ["HEKLA", "HLNET", "HLNODE", "Serie-HEKLA"]
    elif "KARMA" in pattern:
        aliases = ["KARMA_MT", "KARMA", "MT_KARMA"]
    elif "ANIMA" in pattern:
        aliases = ["ANIMA", "ANIMAB"]
    elif "HERCOLA" in pattern:
        aliases = ["HERCOLA", "MT_HERCOLA"]
    elif "TACORA" in pattern:
        aliases = ["TACORA", "TA2000"]
    elif "AURA" in pattern or "NIRVA" in pattern or "MIDAS" in pattern:
        aliases = [pattern.split("_")[0] if "_" in pattern else pattern, pattern]
    elif "VIDOMO" in pattern or "VILLEGGIO" in pattern:
        aliases = ["VIDOMO2K", "VIDOMOBTRX_MT", "VILLEGGIO-NG", "VILLEGGIO", "VIDOMO", "VIBASICB_MT"]
    return aliases


def _prodotti_elmo_espansi(prodotti: list) -> set:
    """Chiavi prodotto ELMO incluse alias e risoluzioni colloquiali."""
    up = set()
    for p in prodotti:
        pu = p.upper()
        up.add(pu)
        up.add(ALIAS_ELMO.get(pu, pu))
        if pu in ALIAS_ELMO.values():
            up.add(pu)
    for k, v in ALIAS_ELMO.items():
        if v in up or k in up:
            up.add(k)
            up.add(v)
    return up


def _filtra_pdf_elmo(matches: list, prodotti: list, domanda_low: str) -> list:
    if not matches:
        return matches
    up = _prodotti_elmo_espansi(prodotti)
    n = lambda p: p.name.lower()

    # Regola generica: per ogni famiglia ELMO riconosciuta, solo PDF coerenti e tecnici
    for key in sorted(up, key=len, reverse=True):
        if key not in ELMO_FILENAME_MARKERS:
            continue
        needle = ELMO_FILENAME_MARKERS[key]
        m = [p for p in matches if needle in n(p) and not _is_elmo_pdf_accessorio(p.name, domanda_low)]
        if key in ("PRX", "PROXIMA", "PRX80", "PRX128", "PRX256", "PRX1024",
                   "PROXIMA80", "PROXIMA128", "PROXIMA256", "PROXIMA1024"):
            mt = [p for p in m if "_mt_mt_serie-proxima" in n(p)]
            if mt:
                return mt
            m = [p for p in m if "tacora" not in n(p)]
        if key in ("VIDOMO", "VIDOMO2K", "VILLEGGIO", "VICOMPACT"):
            mt = [p for p in m if "_mt_" in n(p) or "_mu_" in n(p) or "vidomo2k" in n(p)]
            if mt:
                if "VIDOMO2K" in up:
                    v2k = [p for p in mt if "vidomo2k" in n(p) or "villeggio-ng" in n(p)]
                    if v2k:
                        return v2k
                elif up & {"VIDOMO", "VILLEGGIO", "VICOMPACT"} - {"VIDOMO2K"}:
                    vid = [p for p in mt if "vidomo" in n(p) and "vibasic" not in n(p)]
                    if vid:
                        return vid
                return mt
        if key == "PREGIO500":
            serie = [p for p in matches if "serie-pregio" in n(p) and "_mt_" in n(p)]
            if serie:
                return serie
            tech = [p for p in m if "_mt_" in n(p)]
            if tech:
                return tech
        if key == "PREGIO2000":
            serie = [p for p in matches if "serie-pregio" in n(p) and ("pregio2000" in n(p) or "_mt_" in n(p))]
            if serie:
                return serie
        if key == "KARMA":
            mt = [p for p in m if "_mt_" in n(p)]
            if mt:
                return mt
        if key in ("GATEWAY2K", "GATEWAY"):
            mt = [p for p in m if ("_mt_" in n(p) or n(p).startswith("mt_gateway")) and "gateway2k" in n(p)]
            if mt:
                return mt
        if m:
            return m

    if up & {"SUPERIA", "SPR", "SPR256", "SPR512", "SPR2040"}:
        m = [p for p in matches if "superia" in n(p)]
        if m:
            return m
    if not _domanda_elmo_accessorio(domanda_low):
        centrale = [p for p in matches if not _is_elmo_pdf_accessorio(p.name, domanda_low)]
        if centrale:
            return centrale
    return matches


def _priorita_pdf_elmo(path: Path, prodotti: list | None = None) -> tuple:
    n = path.name.lower()
    if _is_elmo_pdf_accessorio(path.name):
        return (9, 2, n)
    up = {p.upper() for p in (prodotti or [])}
    if up & {"SUPERIA", "SPR"} and "superia" not in n:
        return (8, 2, n)
    if up & {"PRX", "PROXIMA", "PRX80", "PRX128", "PRX256", "PRX1024"}:
        if "proxima" not in n and "serie-proxima" not in n:
            return (8, 2, n)
        if "tacora" in n and "proxima" not in n.replace("tacora", ""):
            return (7, 2, n)
    lang = (
        0 if any(x in n for x in ("_it_", ".it.", "_ita", "italian", "ita_")) else
        2 if any(x in n for x in ("_de_", ".de.", "german", "anweisungen")) else 1
    )
    if up & {"VIDOMO", "VILLEGGIO", "VICOMPACT"} - {"VIDOMO2K"}:
        if "vidomo" in n and "vibasic" not in n and "_mt_" in n:
            return (0, lang, n)
        if "vibasic" in n or "vicompact" in n:
            return (6, lang, n)
    if "PREGIO500" in up:
        if "serie-pregio" in n and "_mt_" in n:
            return (0, lang, n)
        if "pregio500" in n and "_doc" not in n and "_ft_" not in n:
            return (2, lang, n)
    for key in sorted(up, key=len, reverse=True):
        needle = ELMO_FILENAME_MARKERS.get(key)
        if needle and needle in n:
            if "_mt_mt_" in n or (n.startswith("mt_") and needle.replace("-", "") in n.replace("-", "")):
                return (0, lang, n)
            if "_mt_" in n or "_mp_" in n:
                return (1, lang, n)
            if "_mu_" in n:
                return (2, lang, n)
            if "_ft_" in n or "_dc_" in n or "_ce_" in n:
                return (9, lang, n)
    if up & {"GATEWAY2K", "GATEWAY"} or "gateway2k" in n:
        if "_mt_" in n or n.startswith("mt_gateway"):
            return (0, lang, n)
        if "_ce_" in n or "_imq" in n or "_ft_" in n or "_dc_" in n:
            return (9, lang, n)
    if "_mt_mt_serie-proxima" in n:
        return (0, lang, n)
    if "_mt_mt_" in n:
        return (1, lang, n)
    if "elmo_" in n and "_mt_" in n and "serie-proxima" in n:
        return (2, lang, n)
    if "_mp_" in n or "_pm_" in n:
        return (3, lang, n)
    if "_mu_" in n or ("manuale" in n and "utente" in n):
        return (4, lang, n)
    if "installation" in n or "install" in n:
        return (5, lang, n)
    if "user" in n or "manual" in n:
        return (6, lang, n)
    return (7, lang, n)


def _cerca_pdf_elmo_pattern(pattern: str, prodotti: list, domanda_low: str) -> list:
    visti = set()
    matches = []
    aliases = _pattern_aliases_elmo(pattern, prodotti)
    for pat in aliases:
        for p in cerca_file_per_mappa(pat, {"ELMO"}, (".pdf",)):
            if p not in visti:
                visti.add(p)
                matches.append(p)
    if not matches:
        for pat in aliases:
            for p in cerca_file_per_mappa(pat, {"ELMO"}, (".pptx",)):
                if p not in visti:
                    visti.add(p)
                    matches.append(p)
    return _filtra_pdf_elmo(matches, prodotti, domanda_low)


# ── NOTIFIER: schede dep vs manuali programmazione centrale ──
NOTIFIER_PDF_ACCESSORIO_MARKERS = ("cpr-", "cwss-", "vesda", "brochure")


def _domanda_notifier_accessorio(domanda_low: str) -> bool:
    return any(k in domanda_low for k in ("vesda", "faast", "scheda cpr"))


def _is_notifier_pdf_accessorio(nome_file: str, domanda_low: str = "", prodotti: list | None = None) -> bool:
    n = nome_file.lower()
    up = {p.upper() for p in (prodotti or [])}
    if any(m in n for m in NOTIFIER_PDF_ACCESSORIO_MARKERS):
        return True
    # Scheda dep singolo dispositivo: ok per domande su pulsante/rivelatore
    if "dep-ita" in n or "dep-eng" in n:
        if any(k in domanda_low for k in ("pulsante", "w5a", "m5a", "m3a", "fsp", "fdm", "call point", "rottura")):
            return False
        if "FSP" in up or "FDM" in up or any(k in domanda_low for k in ("fsp", "fdm", "fumo", "termico")):
            return False
        if "NFS2" in up or any(k in domanda_low for k in ("nfs", "loop", "centrale notifier", "programm")):
            return True
    return False


def _pattern_aliases_notifier(pattern: str) -> list:
    aliases = [pattern]
    if "AM-8200N" in pattern:
        aliases = ["AM-8200N manu-prog", "AM-8200N", "AM8200N"]
    elif "NFS2" in pattern:
        aliases = ["NFS2Plus", "NFS2", "NFS-2"]
    elif "M5A-W5A" in pattern or "M5A" in pattern:
        aliases = ["M5A-W5A", "Pulsanti Indirizzati", "W5A", "M5A"]
    elif "FSP" in pattern:
        aliases = ["FSP", "fumo"]
    elif "ID3000" in pattern:
        aliases = ["ID3000", "ID-3000"]
    return aliases


def _filtra_pdf_notifier(matches: list, prodotti: list, domanda_low: str) -> list:
    if not matches:
        return matches
    up = {p.upper() for p in prodotti}
    n = lambda p: p.name.lower()
    if "NFS2" in up or any(k in domanda_low for k in ("nfs2", "nfs-2", "loop notifier", "centrale notifier")):
        m = [p for p in matches if "nfs" in n(p) or ("manu-prog" in n(p) and "am-" not in n(p))]
        if m:
            return m
    if up & {"M5A", "W5A", "M3A"} or any(k in domanda_low for k in ("pulsante", "w5a", "m5a", "call point", "rottura")):
        m = [p for p in matches if "m5a" in n(p) or "w5a" in n(p) or "pulsanti" in n(p)]
        if m:
            return m
    if "FSP" in up or "fumo" in domanda_low:
        m = [p for p in matches if "fsp" in n(p)]
        if m:
            return m
    if not _domanda_notifier_accessorio(domanda_low):
        centrale = [p for p in matches if not _is_notifier_pdf_accessorio(p.name, domanda_low, prodotti)]
        if centrale:
            return centrale
    return matches


def _priorita_pdf_notifier(path: Path, prodotti: list | None = None, domanda_low: str = "") -> tuple:
    n = path.name.lower()
    if _is_notifier_pdf_accessorio(path.name, domanda_low, prodotti):
        return (9, n)
    return (
        0 if "manu-prog" in n and ("ita" in n or "dep-ita" not in n) else
        1 if "manu-prog" in n or "manuale di prog" in n else
        2 if "nfs" in n and "plus" in n else
        3 if "manuale" in n or "manual" in n else
        4 if "dep-ita" in n else
        5 if "datasheet" in n or "scheda" in n else 6,
        n,
    )


def _cerca_pdf_notifier_pattern(pattern: str, prodotti: list, domanda_low: str) -> list:
    visti = set()
    matches = []
    for pat in _pattern_aliases_notifier(pattern):
        for p in cerca_file_per_mappa(pat, {"NOTIFIER"}, (".pdf",)):
            if p not in visti:
                visti.add(p)
                matches.append(p)
    return _filtra_pdf_notifier(matches, prodotti, domanda_low)


# ── ESSER: pochi file — filtri per famiglia (FlexES / IQ8 / Winmag) ──
def _is_esser_pdf_accessorio(nome_file: str, domanda_low: str = "") -> bool:
    return False


def _pattern_aliases_esser(pattern: str) -> list:
    if "FlexES" in pattern or "flex" in pattern.lower():
        return ["FlexES", "Centrale FlexES", "flexes"]
    if "esserbus" in pattern.lower() or "iq8" in pattern.lower():
        return [
            "esserbus", "esserbus-PLus", "Rivelazione automatica esserbus",
            "Tecnologia esserbus", "IQ8", "8000",
        ]
    if "Winmag" in pattern:
        return ["Winmag", "WinMag"]
    return [pattern]


def _filtra_pdf_esser(matches: list, domanda_low: str) -> list:
    if not matches:
        return matches
    if any(k in domanda_low for k in ("flexes", "flex es", "centrale flex")):
        m = [p for p in matches if "flexes" in p.name.lower()]
        if m:
            return m
    if any(k in domanda_low for k in ("iq8", "esserbus", "8000", "esser")):
        m = [p for p in matches if "esserbus" in p.name.lower()]
        if m:
            return m
        m = [p for p in matches if "iq8" in p.name.lower()]
        if m:
            return m
    centrale = [p for p in matches if not _is_esser_pdf_accessorio(p.name, domanda_low)]
    return centrale if centrale else matches


def _priorita_pdf_esser(path: Path, domanda_low: str = "") -> tuple:
    n = path.name.lower()
    if _is_esser_pdf_accessorio(path.name, domanda_low):
        return (9, n)
    return (
        0 if "centrale" in n and "flexes" in n else
        1 if "flexes" in n else
        2 if "esserbus" in n or "iq8" in n else
        3 if "winmag" in n else 4,
        n,
    )


def _cerca_pdf_esser_pattern(pattern: str, prodotti: list, domanda_low: str) -> list:
    visti = set()
    matches = []
    for pat in _pattern_aliases_esser(pattern):
        for p in cerca_file_per_mappa(pat, {"ESSER"}, (".pdf",)):
            if p not in visti:
                visti.add(p)
                matches.append(p)
    return _filtra_pdf_esser(matches, domanda_low)


# ── DAHUA: datasheet vs manuale utente/operativo ──
DAHUA_PDF_ACCESSORIO_MARKERS = (
    "_datasheet_", "datasheet_", "_dimensions", "installation_method",
    "brochure", "marketing", "_zip_", "quick_guide",
)


def _domanda_dahua_accessorio(domanda_low: str) -> bool:
    return any(k in domanda_low for k in ("datasheet", "scheda tecnica", "dimensioni"))


def _is_dahua_pdf_accessorio(nome_file: str, domanda_low: str = "") -> bool:
    n = nome_file.lower()
    if _domanda_dahua_accessorio(domanda_low):
        return False
    if any(m in n for m in DAHUA_PDF_ACCESSORIO_MARKERS):
        if "camera-accessories-selection" in n and "datasheet" not in n:
            return False
        if "manuale_utente" in n or "operation_manual" in n or "operation-manual" in n:
            return False
        return True
    if "datasheet" in n:
        return True
    return False


def _pattern_aliases_dahua(pattern: str) -> list:
    if "Camera-Accessories" in pattern:
        return ["Camera-Accessories-Selection", "Accessories-Selection"]
    if "ANPR" in pattern:
        return ["ANPR", "Smart_ANPR", "targa"]
    if "Network_Camera_Web" in pattern or "Operation" in pattern:
        return [
            "Network_Camera_Web_5.0_Operation",
            "Web_5.0_Operation_Manual",
            "Web_3.0_Operation",
            "Operation_Manual",
        ]
    return [pattern]


def _filtra_pdf_dahua(matches: list, domanda_low: str) -> list:
    if not matches:
        return matches
    if not _domanda_dahua_accessorio(domanda_low):
        centrale = [p for p in matches if not _is_dahua_pdf_accessorio(p.name, domanda_low)]
        if centrale:
            return centrale
    return matches


def _priorita_pdf_dahua(path: Path, domanda_low: str = "") -> tuple:
    n = path.name.lower()
    if _is_dahua_pdf_accessorio(path.name, domanda_low):
        return (9, n)
    return (
        0 if "camera-accessories-selection" in n and "datasheet" not in n else
        1 if "camera-accessories-selection" in n else
        2 if "operation_manual" in n or "operation-manual" in n or "web_5.0" in n else
        3 if "manuale_utente" in n or "manuale utente" in n else
        4 if "smart_anpr" in n or "anpr" in n else
        5 if "installation" in n and "method" not in n else
        6 if "datasheet" in n else 7,
        n,
    )


def _cerca_pdf_dahua_pattern(pattern: str, prodotti: list, domanda_low: str) -> list:
    visti = set()
    matches = []
    for pat in _pattern_aliases_dahua(pattern):
        for p in cerca_file_per_mappa(pat, {"DAHUA"}, (".pdf",)):
            if p not in visti:
                visti.add(p)
                matches.append(p)
    return _filtra_pdf_dahua(matches, domanda_low)


def _cerca_pdf_marca_pattern(
    marca: str, pattern: str, prodotti: list, domanda_low: str
) -> list:
    """Ricerca PDF con pattern multipli, filtri accessori e coerenza prodotto."""
    if marca == "RISCO":
        return _cerca_pdf_risco_pattern(pattern, prodotti, domanda_low)
    if marca == "ELMO":
        return _cerca_pdf_elmo_pattern(pattern, prodotti, domanda_low)
    if marca == "NOTIFIER":
        return _cerca_pdf_notifier_pattern(pattern, prodotti, domanda_low)
    if marca == "ESSER":
        return _cerca_pdf_esser_pattern(pattern, prodotti, domanda_low)
    if marca == "DAHUA":
        return _cerca_pdf_dahua_pattern(pattern, prodotti, domanda_low)
    matches = []
    for p in cerca_file_per_mappa(pattern, {marca}, (".pdf",)):
        matches.append(p)
    return matches


def _pattern_manuale_per_marca(marca: str, domanda_low: str, prodotti: list) -> str | None:
    """Sceglie il pattern PDF più adatto per marca, prodotti e testo domanda."""
    prodotti_up = [p.upper() for p in prodotti]
    for keywords, pattern, m, prodotti_req in DOCUMENTO_MANUALE_MAP:
        if m != marca:
            continue
        match_kw = keywords and any(k in domanda_low for k in keywords)
        match_prod = prodotti_req and any(p in prodotti_up for p in prodotti_req)
        if not match_kw and not match_prod:
            continue
        if prodotti_req and not match_prod and not match_kw:
            continue
        return pattern
    return MARCA_FALLBACK_PATTERN.get(marca)


def _priorita_pdf_marca(
    marca: str, path: Path, prodotti: list | None = None, domanda_low: str = ""
) -> tuple:
    """Ordina i candidati PDF: manuali tecnici/programmazione prima di accessori/datasheet."""
    if marca == "RISCO":
        return _priorita_pdf_risco(path, prodotti)
    if marca == "ELMO":
        return _priorita_pdf_elmo(path, prodotti)
    if marca == "NOTIFIER":
        return _priorita_pdf_notifier(path, prodotti, domanda_low)
    if marca == "ESSER":
        return _priorita_pdf_esser(path, domanda_low)
    if marca == "DAHUA":
        return _priorita_pdf_dahua(path, domanda_low)
    return (3, path.name.lower())


def carica_manuali_colloquiali_per_domanda(
    domanda: str, prodotti: list, marche_attese: set | None = None
) -> tuple[str, set, set]:
    """
    Carica manuali mirati per domande in linguaggio naturale (tutte le marche).
    Ritorna (contesto, fonti, marche_caricate).
    """
    domanda_low = domanda.lower()
    marche = _marche_da_domanda_colloquiale(domanda, prodotti, marche_attese)
    if not marche:
        return "", set(), set()

    contesto = ""
    fonti = set()
    marche_caricate = set()

    for marca in sorted(marche):
        operativa = _domanda_operativa_per_marca(domanda_low, marca)
        if not operativa and not prodotti and not _menziona_marca(domanda_low, marca):
            continue

        pattern = _pattern_manuale_per_marca(marca, domanda_low, prodotti)
        if not pattern:
            continue

        matches = _cerca_pdf_marca_pattern(marca, pattern, prodotti, domanda_low)
        if not matches and marca != "DAHUA":
            for p in cerca_file_per_mappa(pattern, {marca}, (".pptx",)):
                matches.append(p)
        if not matches:
            continue

        matches.sort(key=lambda p: _priorita_pdf_marca(marca, p, prodotti, domanda_low))
        pdf = matches[0]
        testo = estrai_testo_file(pdf)
        if not testo:
            continue

        limite = 8000 if marca in ("RISCO", "ELMO", "NOTIFIER") else 6000
        contesto += (
            f"\n=== {pdf.name} (manuale {marca} — linguaggio naturale) ===\n"
            f"{testo[:limite]}\n"
        )
        fonti.add(pdf.name)
        marche_caricate.add(marca)
        logger.info(f"Manuale {marca} mappa caricato: {pdf.name}")

    return contesto, fonti, marche_caricate


def _fonte_elmo_e_manuale_tecnico(nome_fonte: str) -> bool:
    """True se il file ELMO è MT/MP/MU e non CE/FT/DC."""
    n = nome_fonte.lower()
    if _is_elmo_pdf_accessorio(nome_fonte):
        return False
    return "_mt_" in n or n.startswith("mt_") or "_mp_" in n or "_mu_" in n


def codice_in_famiglia_marca(codice: str, marca: str) -> bool:
    """True se il codice è un nome famiglia/colloquiale già coperto dalla mappa marca."""
    cu = codice.upper()
    if marca == "RISCO":
        return cu in FAMIGLIE_RISCO or cu in ALIAS_RISCO
    if marca == "ELMO":
        return cu in FAMIGLIE_ELMO or cu in ALIAS_ELMO or ALIAS_ELMO.get(cu) in FAMIGLIE_ELMO
    return False


def _fonte_notifier_e_manuale_tecnico(nome_fonte: str) -> bool:
    n = nome_fonte.lower()
    if _is_notifier_pdf_accessorio(nome_fonte, "", []):
        return False
    return any(x in n for x in ("manu-prog", "manuale", "nfs", "nfs2", "pulsanti", "id3000"))


def _fonte_dahua_e_manuale_tecnico(nome_fonte: str) -> bool:
    return not _is_dahua_pdf_accessorio(nome_fonte, "")


def _mappa_gia_carica_manuale_tecnico(marca: str, fonti_usate: set, codice: str) -> bool:
    """Evita di saltare trova_e_leggi_pdf se la mappa ha caricato solo CE/datasheet."""
    if not codice_in_famiglia_marca(codice, marca):
        return False
    check = {
        "ELMO": _fonte_elmo_e_manuale_tecnico,
        "RISCO": lambda f: not _is_risco_pdf_accessorio(f),
        "NOTIFIER": _fonte_notifier_e_manuale_tecnico,
        "DAHUA": _fonte_dahua_e_manuale_tecnico,
        "ESSER": lambda f: "esserbus" in f.lower() or "flexes" in f.lower() or "centrale" in f.lower(),
    }
    fn = check.get(marca)
    if not fn:
        return False
    for fonte in fonti_usate:
        if (marca_da_nome_file(fonte) or "") != marca:
            continue
        if fn(fonte):
            return True
    return False


async def estrai_codici_dalla_storia(domanda: str, storia: list) -> list:
    """
    Estrae codici modello dalla domanda corrente e dai messaggi precedenti.
    Utile per domande di follow-up tipo "approfondisci" o "dimmi di più".
    Gestisce anche codici con spazio (es. "gateway 2k" -> "GATEWAY2K")
    e codici corti (es. "prx80", "spr256").
    """
    import re

    # Pattern 1: codici standard con trattino (es. IPC-HFW2241T-ZS)
    # Pattern 2: codici alfanumerici con numeri (es. NVR5208EI2, GATEWAY2K, PRX80)
    # Pattern 3: codici corti tipo PRX80, SPR256 (2+ lettere + 2+ numeri)
    pattern_standard = r'[A-Z0-9]{3,}(?:[-][A-Z0-9]+)+|[A-Z]{3,}[0-9]+[A-Z0-9]*|[A-Z]{2,}[0-9]{2,}[A-Z0-9]*'

    domanda_up = domanda.upper()

    # Normalizza "lightsys+" -> "LIGHTSYSPLUS", "lightsys 2" -> "LIGHTSYS2" ecc.
    domanda_up = domanda_up.replace("LIGHTSYS+", "LIGHTSYSPLUS")
    domanda_up = domanda_up.replace("LIGHTSYS PLUS", "LIGHTSYSPLUS")
    domanda_up = domanda_up.replace("LIGHTSYS AIR", "LIGHTSYSAIR")
    domanda_up = domanda_up.replace("WICOMM PRO", "WICOMMPRO")
    domanda_up = domanda_up.replace("AGILITY 4", "AGILITY4")

    # Aggiungi anche nomi colloquiali puri (solo lettere, es. LIGHTSYS, AGILITY)
    pattern_nomi = r'[A-Z]{5,}'  # parole lunghe 5+ lettere potrebbero essere nomi prodotto
    nomi_trovati = re.findall(pattern_nomi, domanda_up)

    # Normalizza spazi in certi pattern comuni (es. "GATEWAY 2K" -> "GATEWAY2K")
    domanda_norm = re.sub(r'([A-Z]+)\s+([0-9]+[A-Z]*)', r'\1\2', domanda_up)

    codici = list(set(
        re.findall(pattern_standard, domanda_up) +
        re.findall(pattern_standard, domanda_norm) +
        nomi_trovati
    ))

    # Filtra codici troppo generici
    ESCLUDI = {"THE", "AND", "FOR", "CON", "PER", "DAL", "DEL", "UNA", "UNO",
                "NON", "CHE", "SUI", "SUL", "ALL", "SIA", "VIA", "ORA", "GHZ",
                "POI", "BUS", "LED", "LAN", "USB", "SIM", "GSM", "LTE", "POE",
                "TCP", "UDP", "FTP", "SSH", "VPN", "DNS", "NTP", "MAC", "RAM",
                "ROM", "CPU", "GPU", "API", "SDK", "CGI", "URL", "HTTP", "HTTPS",
                "PDF", "ZIP", "JPG", "PNG", "AVI", "MP4", "FPS", "HDR", "SDR",
                # Parole comuni italiane che non sono codici modello
                "MANUALE", "MANUALI", "CENTRALE", "CENTRALI", "SENSORE", "SENSORI",
                "RIVELATORE", "TASTIERA", "SIRENA", "TELECAMERA", "TELECAMERE",
                "ZONA", "ZONE", "CODICE", "CODICI", "UTENTE", "UTENTI",
                "INSTALLAZIONE", "PROGRAMMAZIONE", "CONFIGURAZIONE", "PROCEDURA",
                "COME", "COSA", "DOVE", "QUANDO", "QUALE", "QUALI", "QUESTO",
                "QUESTA", "QUELLO", "QUELLA", "VOGLIO", "VORREI", "POSSO",
                "PUOI", "RIESCI", "FORNISCI", "DIMMI", "MOSTRAMI", "DAMMI",
                "SERIE", "SISTEMA", "IMPIANTO", "ALLARME", "SEGNALE", "USCITA",
                "INGRESSO", "PARTIZIONE", "AREA", "GRUPPO", "RETE", "CLOUD",
                "RISCO", "DAHUA", "ELMO", "NOTIFIER", "ESSER", "SAMSUNG", "BOSCH", "HONEYWELL",
                "DELLA", "DELLO", "DEGLI", "DELLE", "FORNISCIMI", "CONFIGURO",
                "DIMMI", "MOSTRAMI", "CERCAMI", "TROVAMI", "INSTALLARE",
                "COLLEGARE", "CONNETTERE", "RESETTARE", "PROGRAMMARE",
                "IMPOSTARE", "ABILITARE", "DISABILITARE", "VERIFICARE",
                "CONTROLLARE", "GESTIRE", "AGGIORNARE", "SCARICARE",
                "QUALI", "QUALE", "QUANTO", "QUANTA", "QUANTI", "QUANTE",
                "TUTTO", "TUTTI", "TUTTA", "TUTTE", "PRIMO", "PRIMA",
                "SECONDO", "SECONDA", "NUOVO", "NUOVA", "VECCHIO",
                "GRANDE", "PICCOLO", "ESTERNO", "INTERNO", "WIRELESS",
                "CABLATO", "DIGITALE", "ANALOGICO", "INTEGRATO",
                "VERSIONE", "MODELLO", "PRODOTTO", "DISPOSITIVO",
                "FORNICA", "FORNICIMI", "FORNIRE", "INVIARMI",
                "MANDAMI", "INVIAMI", "CERCAMI", "AIUTAMI", "SPIEGAMI",
                "CABLATO", "FUNZIONE", "FUNZIONA", "FUNZIONAMENTO",
                # Parole che matchano per errore nomi di file (es. COMANDI -> TELECOMANDI)
                "COMANDI", "COMANDO", "INVIARE", "INVIANDO", "POSSIBILITA",
                "HANNO", "HANNO", "POSSO", "POSSO", "POSSONO", "POSSIBILE",
                "INVIO", "INVII", "INVIO", "TRAMITE", "MEDIANTE", "ATTRAVERSO",
                "NOTIFICA", "NOTIFICHE", "EVENTO", "EVENTI", "AZIONE", "AZIONI",
                "ALLARMI", "ALLARME", "REGOLA", "REGOLE", "TRIGGER", "PORTA",
                "PORTE", "INDIRIZZO", "INDIRIZZI", "RISPOSTA", "RICHIESTA",
                "RICHIESTE", "RISPONDE", "CHIAMATA", "CHIAMATE", "METODO",
                "METODI", "PROTOCOLLO", "PROTOCOLLI", "INTEGRAZIONE", "ACCESSO",
                "POSSIBILITA", "FUNZIONALITA", "SUPPORTO", "SUPPORTA"}
    codici = [c for c in codici if c not in ESCLUDI and codice_modello_plausibile(c)]

    # Se non ci sono codici nella domanda corrente, cerca nella storia
    if not codici and storia:
        for msg in reversed(storia[:-1]):
            testo = msg.get("content", "")
            testo_up = testo.upper()
            testo_norm = re.sub(r'([A-Z]+)\s+([0-9]+[A-Z]*)', r'\1\2', testo_up)
            codici_trovati = list(set(
                re.findall(pattern_standard, testo_up) +
                re.findall(pattern_standard, testo_norm)
            ))
            codici_trovati = [c for c in codici_trovati if c not in ESCLUDI and codice_modello_plausibile(c)]
            if codici_trovati:
                codici = codici_trovati
                logger.info(f"Codici trovati dalla storia: {codici}")
                break

    for prodotto in estrai_prodotti_da_testo(domanda):
        chiave = risolvi_prodotto_colloquiale(prodotto)
        if chiave not in codici:
            codici.append(chiave)

    return codici

MARCHE_FILE = ("DAHUA", "ESSER", "ELMO", "RISCO", "NOTIFIER")


def utente_e_admin(user_id: int) -> bool:
    """True solo se ADMIN_CHAT_ID è configurato e coincide con l'utente."""
    return bool(ADMIN_CHAT_ID) and user_id == ADMIN_CHAT_ID


def path_manuale_sicuro(nome_file: str) -> Path | None:
    """
    Risolve un PDF solo dentro MANUALI_DIR (anti path traversal).
    Accetta solo il nome file, non percorsi relativi.
    """
    if not nome_file or not str(nome_file).strip():
        return None
    nome_pulito = Path(nome_file).name
    if not nome_pulito.lower().endswith(".pdf"):
        return None
    try:
        if hasattr(MANUALI_DIR, "find_file"):
            found = MANUALI_DIR.find_file(nome_pulito)  # type: ignore[attr-defined]
            if found:
                return found
        base = MANUALI_DIR.resolve()
        candidato = (base / nome_pulito).resolve()
        if not str(candidato).startswith(str(base)):
            logger.warning(f"Path PDF rifiutato (fuori da manuali): {nome_file}")
            return None
        if candidato.is_file():
            return candidato
    except OSError as e:
        logger.warning(f"Path PDF non valido '{nome_file}': {e}")
    return None


def _processo_vivo(pid: int) -> bool:
    if pid <= 0:
        return False
    if os.name == "nt":
        import ctypes
        PROCESS_QUERY_LIMITED_INFORMATION = 0x1000
        STILL_ACTIVE = 259
        handle = ctypes.windll.kernel32.OpenProcess(PROCESS_QUERY_LIMITED_INFORMATION, False, pid)
        if not handle:
            return False
        try:
            code = ctypes.c_ulong()
            if ctypes.windll.kernel32.GetExitCodeProcess(handle, ctypes.byref(code)):
                return int(code.value) == STILL_ACTIVE
            return False
        finally:
            ctypes.windll.kernel32.CloseHandle(handle)
    try:
        os.kill(pid, 0)
    except OSError:
        return False
    return True


def verifica_istanza_unica():
    """Registra PID processo (lock già acquisito all'avvio script)."""
    PID_FILE.write_text(str(os.getpid()), encoding="utf-8")


def rimuovi_pid_file():
    global _early_lock_fd, _instance_lock_fd
    try:
        if PID_FILE.exists() and PID_FILE.read_text(encoding="utf-8").strip() == str(os.getpid()):
            PID_FILE.unlink()
    except Exception:
        pass
    if _early_lock_fd is not None:
        try:
            os.close(_early_lock_fd)
        except Exception:
            pass
        try:
            if _EARLY_LOCK_FILE.exists():
                _EARLY_LOCK_FILE.unlink(missing_ok=True)
        except Exception:
            pass
    _early_lock_fd = None
    _instance_lock_fd = None


def codice_modello_plausibile(codice: str) -> bool:
    """True solo per stringhe che assomigliano a un codice prodotto, non parole italiane."""
    c = codice.upper().strip()
    if not c or c in PAROLE_NON_CODICI or len(c) < 4:
        return False
    if marca_da_codice_modello(c):
        return True
    if any(ch.isdigit() for ch in c):
        return len(c) >= 5
    if "-" in c and len(c) >= 6:
        return True
    return False

# Prefissi codice modello → marca (priorità su parole generiche tipo "staffa")
PREFISSI_CODICE_MARCA = (
    ("DAHUA", (
        "IPC-H", "IPC-", "HDW", "HFW", "HDB", "HAC", "HAD", "HMW", "HFE", "HFS",
        "NVR", "XVR", "DHI-", "DHI", "TPC-", "TPC", "PTZ", "PSDW", "PSSD",
        "SD1", "SD2", "SD3", "SD4", "SD5", "SD6", "SD7", "SD8", "SD9",
    )),
    ("ELMO", (
        "PRX", "SPR", "PREGIO", "KARMA", "AURA", "ANIMA", "GATEWAY2K", "GATEWAY",
        "VIDOMO", "VILLEGGIO", "HEKLA", "HERCOLA", "TACORA", "NIRVA", "MIDAS",
    )),
    ("RISCO", ("LIGHTSYS", "WICOMM", "AGILITY", "RP", "GT")),
    ("NOTIFIER", ("AM-", "AM8", "NFS", "NFB", "NFW", "FSP", "FDM", "M3A", "M5A", "W5A", "ID3000")),
    ("ESSER", ("IQ8", "FX2", "FX10", "FX18", "8000C", "9200")),
)


def marca_da_nome_file(nome: str) -> str | None:
    """Estrae la marca dal prefisso del nome file (es. DAHUA_..., ESSER_...)."""
    testa = nome.split("_")[0].upper()
    return testa if testa in MARCHE_FILE else None


def marca_da_codice_modello(codice: str) -> str | None:
    """Rileva la marca da un codice prodotto (es. HFW3849T-ZS-IL → DAHUA)."""
    c = codice.upper()
    for marca, prefissi in PREFISSI_CODICE_MARCA:
        for p in prefissi:
            if c.startswith(p):
                return marca
    return None


def rileva_marche_attese(domanda: str, codici: list) -> set:
    """
    Marche rilevanti per la domanda.
    I codici modello hanno priorità sulle parole generiche (staffa, montaggio, ecc.).
    """
    marche = set()
    for codice in codici:
        m = marca_da_codice_modello(codice)
        if m:
            marche.add(m)
    if marche:
        return marche
    m = _rileva_marca_domanda(domanda)
    return {m} if m else set()


def filtra_fonti_per_marca(fonti: set, marche: set | None) -> list:
    """Mostra in 'Fonti' solo i file coerenti con la marca della domanda."""
    if not marche:
        return sorted(fonti)
    risultato = []
    for nome in sorted(fonti):
        if nome.startswith("linee_guida") or nome.startswith("ricerca_prefisso"):
            risultato.append(nome)
            continue
        marca_file = marca_da_nome_file(nome)
        if marca_file and marca_file in marche:
            risultato.append(nome)
    return risultato


def cerca_file_per_mappa(pattern: str, marche: set, estensioni: tuple = (".pdf",)) -> list:
    """Cerca manuali per pattern limitando alle marche consentite."""
    trovati = []
    for marca in sorted(marche):
        for ext in estensioni:
            trovati.extend(MANUALI_DIR.glob(f"**/*{marca}*{pattern}*{ext}"))
    visti = set()
    unici = []
    for p in trovati:
        if p not in visti:
            visti.add(p)
            unici.append(p)
    return unici


def _rileva_marca_domanda(testo: str) -> str | None:
    """Rileva la marca da una stringa di testo (domanda o storia)."""
    t = testo.upper()
    if any(x in t for x in ["DAHUA", "IPC-", "IPC ", "NVR", "HDW", "HFW", "HDB", "XVR", "DHI-"]):
        return "DAHUA"
    if any(x in t for x in ["RISCO", "LIGHTSYS", "AGILITY", "WICOMM"]):
        return "RISCO"
    if any(x in t for x in ["ESSER", "IQ8", "IQ8CONTROL", "ESSERBUS",
                              "IQ8QUAD", "IQ8ALARM", "TOOL8000",
                              "FLEXES", "FLEX ES", "WINMAG", "ESSERNET",
                              "CENTRALE FLEX", "FLEX CONTROL"]):
        return "ESSER"
    # "FLEX" solo come parola intera (evita falsi positivi su altre parole)
    if re.search(r"\bFLEX\b", t):
        return "ESSER"
    if any(x in t for x in ["ELMO", "PROXIMA", "PRX", "PREGIO", "KARMA", "ANIMA",
                              "TACORA", "SUPERIA", "TRIAL", "NIRVA", "VIDOMO"]):
        return "ELMO"
    if any(x in t for x in ["NOTIFIER", "AM-8200", "AM-8100", "NFS", "W5A", "M5A",
                              "AM2000", "AM6000", "AM1000", "AM4000", "AM82",
                              "NFXI", "FAAST", "VESDA", "LCD600"]):
        return "NOTIFIER"
    return None

def storia_testo_breve(storia: list) -> str:
    """Restituisce il testo degli ultimi 4 messaggi della storia."""
    if not storia:
        return ""
    return " ".join(m.get("content", "")[:200] for m in storia[-4:])

async def genera_risposta(domanda: str, storia: list = None) -> dict:
    """
    Pipeline completa:
    1. Se c'è un codice modello nella domanda (o nella storia), legge direttamente il PDF
    2. Integra sempre con ricerca semantica per info da altri manuali
    3. Chiama Claude con il contesto recuperato
    4. Genera audio con edge-tts
    """
    import re
    client_ant = Anthropic(api_key=ANTHROPIC_KEY)
    codici = await estrai_codici_dalla_storia(domanda, storia or [])
    richiesta_pdf = is_richiesta_pdf(domanda)
    marche_attese = rileva_marche_attese(domanda, codici)
    if marche_attese:
        logger.info(f"Marche attese (priorità codici modello): {marche_attese}")

    contesto = ""
    fonti_usate = set()
    pdf_paths = []  # percorsi file PDF da allegare

    # Se è una richiesta esplicita di PDF senza codice modello,
    # cerca i PDF già menzionati nella storia recente
    if richiesta_pdf and not codici and storia:
        for msg in reversed(storia[:-1]):
            testo_msg = msg.get("content", "")
            # Cerca nomi file PDF nella storia (risposte precedenti)
            import re as _re
            nomi_trovati = _re.findall(r'[\w\-]+\.pdf', testo_msg, _re.IGNORECASE)
            if nomi_trovati:
                for nome in nomi_trovati[:3]:
                    p = path_manuale_sicuro(nome)
                    if p:
                        pdf_paths.append(p)
                        fonti_usate.add(p.name)
                if pdf_paths:
                    logger.info(f"PDF dalla storia: {[p.name for p in pdf_paths]}")
                    break

    # 0. Se è una domanda di dimensionamento, carica le linee guida come primo contesto
    # Controlla anche la storia recente (per messaggi di follow-up)
    storia_testo = " ".join(m.get("content", "") for m in (storia[-4:] if storia else []))
    prodotti_colloquiali = estrai_prodotti_da_testo(domanda)
    manuali_mappa_caricati = set()

    is_domanda_dim = (
        is_domanda_dimensionamento(domanda, codici, marche_attese)
        or is_domanda_dimensionamento(storia_testo, [], marche_attese)
    )
    ha_datasheet_dimensionamento = False

    if is_domanda_dim:
        linee_guida = carica_linee_guida()
        if linee_guida:
            sezioni_lg = seleziona_sezioni_linee_guida(domanda, linee_guida)
            contesto += (
                "\n=== LINEE GUIDA DIMENSIONAMENTO — CATALOGO MODELLI AMMESSI ===\n"
                f"{sezioni_lg}\n"
            )
            fonti_usate.add("linee_guida_dimensionamento.txt")
            logger.info("Linee guida (sezioni pertinenti) caricate per dimensionamento")

            ctx_datasheet, fonti_datasheet = carica_datasheet_per_dimensionamento(
                domanda, linee_guida
            )
            if ctx_datasheet:
                contesto += (
                    "\n=== DATASHEET MODELLI CANDIDATI — VERIFICA SPECIFICHE TECNICHE ===\n"
                    f"{ctx_datasheet}\n"
                )
                fonti_usate.update(fonti_datasheet)
                ha_datasheet_dimensionamento = True
                logger.info(
                    f"Datasheet dimensionamento: {len(fonti_datasheet)} file per candidati linee guida"
                )

    # 0b. Manuali mirati per linguaggio naturale (RISCO, ELMO, NOTIFIER, ESSER, DAHUA)
    ctx_map, fonti_map, marche_map = carica_manuali_colloquiali_per_domanda(
        domanda, prodotti_colloquiali or codici, marche_attese
    )
    if ctx_map:
        contesto += ctx_map
        fonti_usate.update(fonti_map)
        manuali_mappa_caricati = marche_map
        if not marche_attese:
            marche_attese = marche_map

    # 1. Lettura diretta del PDF se c'è un codice modello nella domanda
    for codice in codici:
        if not codice_modello_plausibile(codice):
            continue
        if _mappa_gia_carica_manuale_tecnico("ELMO", fonti_usate, codice) or any(
            _mappa_gia_carica_manuale_tecnico(m, fonti_usate, codice)
            for m in manuali_mappa_caricati
        ):
            continue
        testo_pdf, nomi_file = trova_e_leggi_pdf(codice)
        if testo_pdf:
            for nome in nomi_file:
                fonti_usate.add(nome)
                # Salva il path completo per allegarlo in chat
                p = path_manuale_sicuro(nome)
                if p:
                    pdf_paths.append(p)
            contesto += testo_pdf
            logger.info(f"Contesto da {len(nomi_file)} PDF diretti per '{codice}'")
            break  # un codice modello alla volta
        else:
            # Codice non trovato direttamente — cerca modelli simili con quel prefisso
            modelli_simili = trova_modelli_simili(codice)
            if modelli_simili:
                lista = ", ".join(modelli_simili)
                contesto += (
                    f"\n=== MODELLI TROVATI IN ARCHIVIO CON PREFISSO '{codice}' ===\n"
                    f"Non ho trovato un modello esatto per '{codice}'. "
                    f"In archivio ho documentazione per questi modelli simili: {lista}.\n"
                    f"Chiedi all'utente di specificare il modello esatto tra quelli disponibili "
                    f"o fornisci informazioni basandoti sulle linee guida staffaggi.\n"
                )
                fonti_usate.add(f"ricerca_prefisso_{codice}")
                logger.info(f"Modelli simili per '{codice}': {modelli_simili}")

            # Fallback Dahua: se il codice è chiaramente Dahua ma non ha un PDF dedicato,
            # carica il manuale operativo Web generico (Web_5.0 o Web_3.0) se presente
            codice_up_fb = codice.upper()
            PREFISSI_DAHUA_FB = ("HDW","HFW","HDB","HAC","HAD","HMW","IPC","NVR","XVR",
                                  "DHI","TPC","PTZ","SD1","SD2","SD3","SD4","SD5","SD6",
                                  "SD7","SD8","SD9","HFE","HFS")
            is_dahua_codice = any(codice_up_fb.startswith(p) or codice_up_fb.startswith("IPC-"+p)
                                  for p in PREFISSI_DAHUA_FB)
            if is_dahua_codice and not contesto.strip():
                MANUALI_GENERICI_DAHUA = [
                    "Dahua_Network_Camera_Web_5.0_Operation_Manual",
                    "Dahua_Network_Camera_Web_3.0_Operation_Manual",
                    "Network_Camera_Web_5.0_Operation_Manual",
                    "Network_Camera_Web_3.0_Operation_Manual",
                ]
                for pattern in MANUALI_GENERICI_DAHUA:
                    matches = list(MANUALI_DIR.glob(f"**/*{pattern}*.pdf"))
                    if not matches:
                        matches = list(MANUALI_DIR.glob(f"**/*Web*5.0*Operation*Manual*.pdf"))
                    if not matches:
                        matches = list(MANUALI_DIR.glob(f"**/*Web*3.0*Operation*Manual*.pdf"))
                    if matches:
                        pdf_gen = matches[0]
                        testo_gen = estrai_testo_file(pdf_gen)
                        if testo_gen:
                            contesto += f"\n=== {pdf_gen.name} (manuale operativo generico) ===\n{testo_gen[:8000]}\n"
                            fonti_usate.add(pdf_gen.name)
                            # NON aggiunge a pdf_paths: il PDF viene allegato solo se richiesto esplicitamente
                            logger.info(f"[DIAG3] Fallback manuale generico Dahua: {pdf_gen.name}")
                        break

    # 2. Ricerca semantica — con filtro marca e soglia score alta
    #
    # Regole:
    # - Se la domanda menziona una marca specifica → solo quella marca
    # - Se abbiamo già trovato PDF diretti → solo la stessa marca dei PDF trovati
    # - Se abbiamo solo linee_guida → usa marca_attesa se definita
    # - Soglia minima 0.65 per evitare falsi positivi
    # - Mai aggiungere chunk se già abbiamo abbastanza contesto (>8000 chars)

    SOGLIA_SEMANTICA = 0.65
    MAX_CONTESTO = 22000 if is_domanda_dim else 8000

    marca_attesa = marche_attese if marche_attese else None

    # Determina marche già presenti nei PDF trovati (escludi linee_guida e ricerche)
    marche_dai_pdf = set()
    for fonte in fonti_usate:
        m = marca_da_nome_file(fonte)
        if m:
            marche_dai_pdf.add(m)

    # Guardia anti-contaminazione: se marca_attesa è già definita e marche_dai_pdf
    # contiene una marca diversa, ignora marche_dai_pdf (era un match spurio su
    # una parola generica come COMANDI -> TELECOMANDI RISCO)
    if marca_attesa and marche_dai_pdf and not marche_dai_pdf.issubset(marca_attesa):
        logger.info(f"[DIAG2] marche_dai_pdf={marche_dai_pdf} scartate: conflitto con marca_attesa={marca_attesa}")
        marche_dai_pdf = set()

    # Se non abbiamo ancora una marca, cerca nella storia della conversazione
    # (gestisce follow-up tipo "allora verifica nel manuale della proxima")
    marche_dalla_storia = set()
    if not marca_attesa and not marche_dai_pdf and storia:
        KEYWORD_MARCA_STORIA = {
            "ELMO":     {"ELMO", "PROXIMA", "PRX", "SUPERIA", "SPR", "PREGIO", "KARMA",
                         "AURA", "NIRVA", "VIDOMO", "VILLEGGIO", "GATEWAY2K", "TACORA",
                         "HEKLA", "HERCOLA", "ECONNECT", "BROWSERONE"},
            "DAHUA":    {"DAHUA", "NVR", "IPC-", "ANPR", "VARCO", "TIOC", "WIZSENSE",
                         "HDW", "HFW", "HDB", "HAC", "HAD", "TPC-", "XVR", "DHI-",
                         "WIZMIND", "WIZSEEK", "ACUPICK", "FULLCOLOR", "STARLIGHT"},
            "RISCO":    {"RISCO", "LIGHTSYS", "WICOMM", "AGILITY"},
            "NOTIFIER": {"NOTIFIER", "NFS", "NFB", "AM-8200", "AM8200", "AM-8200N",
                         "AM-8200G", "AM2000", "AM6000", "AM1000", "AM4000", "AM82",
                         "NFS2", "NCD", "FCM", "FDM", "FMM", "FZM",
                         "FSP", "FST", "M3A", "M5A", "W5A", "ID3000", "ID50", "ONYX"},
            "ESSER":    {"ESSER", "IQ8", "IQ8CONTROL", "ESSERBUS", "8000C", "9200",
                         "IQ8QUAD", "IQ8ALARM", "TOOL8000", "FAAST", "IQ8TAL",
                         "FLEXES", "FLEX", "WINMAG", "ESSERNET"},
        }
        for msg in storia[-6:]:
            testo_up = msg.get("content", "").upper()
            for marca, keywords in KEYWORD_MARCA_STORIA.items():
                if any(kw in testo_up for kw in keywords):
                    marche_dalla_storia.add(marca)
                    break
        if marche_dalla_storia:
            logger.info(f"Marca ereditata dalla storia: {marche_dalla_storia}")

    # Priorità: marca_attesa > marche_dai_pdf > marche_dalla_storia > nessun filtro
    filtro_marca = marca_attesa or marche_dai_pdf or marche_dalla_storia or None
    logger.info(f"[DIAG2] filtro_marca finale: {filtro_marca} (marca_attesa={marca_attesa}, dai_pdf={marche_dai_pdf}, dalla_storia={marche_dalla_storia})")

    # Salta ricerca semantica se abbiamo già documentazione mirata per la marca
    ha_doc_mirato = bool(marche_attese) and any(
        marca_da_nome_file(f) in marche_attese for f in fonti_usate
    )
    if ha_datasheet_dimensionamento:
        ha_doc_mirato = True

    # Non aggiungere semantica se il contesto è già abbondante o c'è già un doc mirato
    if len(contesto) < MAX_CONTESTO and not ha_doc_mirato:
        chunks = cerca_nei_manuali(domanda)
        logger.info(f"[DIAG2] chunks ricevuti: {len(chunks)}, filtro_marca={filtro_marca}")
        for chunk in chunks:
            fonte = chunk["fonte"]
            fonte_marca = marca_da_nome_file(fonte) or ""
            # Applica filtro marca
            if filtro_marca and fonte_marca not in filtro_marca:
                logger.info(f"[DIAG2] ESCLUSO (marca={fonte_marca}): {fonte[:60]}")
                continue
            # Applica soglia score
            if chunk["score"] < SOGLIA_SEMANTICA:
                logger.info(f"[DIAG2] ESCLUSO (score={chunk['score']}<{SOGLIA_SEMANTICA}): {fonte[:60]}")
                continue
            if fonte not in fonti_usate:
                contesto += f"\n--- Estratto da '{fonte}' ---\n"
                contesto += chunk["testo"] + "\n"
                fonti_usate.add(fonte)
                logger.info(f"[DIAG2] AGGIUNTO: {fonte[:60]}")
            # Ferma se il contesto è sufficiente
            if len(contesto) >= MAX_CONTESTO:
                break

    # Se ha trovato PDF dalla storia per una richiesta esplicita, restituiscili subito
    if richiesta_pdf and pdf_paths and not contesto.strip():
        return {
            "testo": f"Ecco {'il PDF' if len(pdf_paths)==1 else 'i PDF'} richiesti.",
            "audio_path": None,
            "fonti": filtra_fonti_per_marca(fonti_usate, filtro_marca or marca_attesa),
            "pdf_paths": pdf_paths,
            "pdf_richiesti": True,
            "trovato": True,
        }

    if not contesto.strip():
        # Messaggio più specifico se avevamo rilevato una marca
        if marca_attesa == {"DAHUA"}:
            testo_nf = (
                "Non ho trovato il manuale di questo modello Dahua in archivio. "
                "Posso cercarlo online automaticamente: invia il codice esatto del modello "
                "(es. IPC-HDW3849ZS-IL) e proverò a scaricare la documentazione dal sito Dahua."
            )
        elif marca_attesa == {"ELMO"}:
            testo_nf = "Non ho trovato documentazione per questo prodotto ELMO. Verifica di aver caricato il manuale corretto con /aggiorna."
        elif marca_attesa == {"RISCO"}:
            testo_nf = "Non ho trovato documentazione per questo prodotto RISCO. Verifica di aver caricato il manuale corretto con /aggiorna."
        elif marca_attesa == {"NOTIFIER"}:
            testo_nf = "Non ho trovato documentazione per questo prodotto Notifier. Verifica di aver scaricato i manuali con il downloader e indicizzato con /aggiorna."
        elif marca_attesa == {"ESSER"}:
            testo_nf = "Non ho trovato documentazione per questo prodotto ESSER. Verifica di aver caricato i manuali con prefisso ESSER_ e indicizzato con /aggiorna."
        else:
            testo_nf = "Non ho trovato informazioni sufficienti nei manuali per rispondere a questa domanda. Prova a riformularla o verifica che i manuali relativi siano stati caricati."
        return {
            "testo": testo_nf,
            "audio_path": None,
            "fonti": [],
            "pdf_paths": [],
            "trovato": False,
        }

    # 3. Risposta Claude
    try:
        risposta = client_ant.messages.create(
            model="claude-opus-4-5",
            max_tokens=2000,
            system=(
                "Sei un tecnico commerciale specializzato in impianti di videosorveglianza "
                "e antintrusione per la società MAC SYSTEM. "
                "Hai accesso ai manuali tecnici e alle linee guida di dimensionamento. "
                "Rispondi sempre in italiano in modo completo e dettagliato. "
                "\n\n"
                "REGOLA PER IL DIMENSIONAMENTO IMPIANTI: "
                "Quando un utente chiede di dimensionare un impianto ma non ha fornito "
                "tutte le informazioni necessarie, rispondi con UN UNICO messaggio che contiene "
                "tutte le domande necessarie elencate chiaramente, e invita l'utente a rispondere "
                "in un unico messaggio con tutte le informazioni. "
                "Le informazioni necessarie sono: "
                "1) tipo di sito (casa, negozio, capannone, ufficio...), "
                "2) numero di telecamere o zone da coprire, "
                "3) installazione interna/esterna/mista, "
                "4) deterrenza attiva richiesta si/no, "
                "5) giorni di registrazione desiderati. "
                "Quando invece l'utente ha già fornito tutte queste informazioni in un unico messaggio, "
                "proponi direttamente la soluzione completa con i modelli specifici. "
                "WORKFLOW SELEZIONE PRODOTTI (linee guida + datasheet): "
                "1) Identifica i candidati SOLO tra i codici modello in linee_guida_dimensionamento.txt. "
                "2) Per ogni caratteristica richiesta (audio bidirezionale, IP, WDR, risoluzione, ecc.) "
                "VERIFICA nei datasheet forniti nel contesto. "
                "3) Proponi SOLO modelli presenti nelle linee guida E confermati dal datasheet "
                "per la caratteristica richiesta. "
                "4) Escludi candidati delle linee guida se il datasheet non conferma la caratteristica. "
                "5) NON proporre modelli assenti dalle linee guida, anche se compaiono nei datasheet. "
                "6) Se un candidato non ha datasheet in archivio, indicalo senza inventare specifiche. "
                "Audio bidirezionale significa microfono E altoparlante/speaker integrati per comunicazione two-way. "
                "Se per una categoria non trovi un modello adatto nell'elenco linee guida, "
                "segnalalo esplicitamente invece di inventare alternative. "
                "PER DOMANDE TECNICHE SUI PRODOTTI: "
                "Cerca SEMPRE attentamente nel testo dei manuali forniti prima di rispondere. "
                "Le specifiche tecniche si trovano nelle sezioni 'Environment', 'Power', "
                "'Structure', 'Technical Specification'. "
                "Se trovi il dato richiesto, riportalo con i valori esatti. "
                "Se l'informazione NON è nei testi forniti, dillo chiaramente senza inventare. "
                "MAI inventare numeri di pagina, sezioni o procedure non presenti nel testo fornito. "
                "Se l'utente chiede 'a che pagina' o 'dove trovo' una informazione, "
                "rispondi solo se il numero di pagina è esplicitamente indicato nel testo fornito, "
                "altrimenti indica la sezione o il capitolo pertinente senza inventare pagine. "
                "Le linee guida di dimensionamento contengono tabelle di compatibilità ufficiali — "
                "quando trovi la risposta nelle linee guida, forniscila con sicurezza senza aggiungere "
                "disclaimer come 'ti consiglio di verificare' o 'non posso confermare'. "
                "Scrivi in testo continuo senza asterischi o formattazione markdown. "
                "NOTA IMPORTANTE: la centrale FLEX o FlexES o FlexES Control "
                "è un prodotto ESSER by Honeywell, NON di EL.MO. "
                "Se l'utente chiede della centrale FLEX senza specificare il produttore, "
                "rispondere sempre con informazioni ESSER FlexES."
            ),
            messages=(
                [m for m in (storia[:-1] if storia and len(storia) > 1 else [])
                 if m["role"] in ("user", "assistant")]
                + [{
                    "role": "user",
                    "content": (
                        f"Domanda: {domanda}\n\n"
                        f"Testo dei manuali:\n{contesto}\n\n"
                        f"Rispondi alla domanda basandoti sul testo fornito."
                    )
                }]
            )
        )
        testo_risposta = risposta.content[0].text
    except Exception as e:
        logger.error(f"Errore Claude: {e}")
        testo_risposta = "Si è verificato un errore durante l'elaborazione della risposta. Riprova tra qualche secondo."

    # 4. Genera audio (solo bot Telegram standalone)
    audio_path = None
    if edge_tts is not None and __name__ == "__main__":
        testo_audio = pulisci_testo(testo_risposta)
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            audio_path = f"risposta_{timestamp}.mp3"
            await edge_tts.Communicate(testo_audio, VOICE).save(audio_path)
            logger.info(f"✅ Audio generato: {audio_path}")
        except Exception as e:
            logger.error(f"Errore TTS: {e}")
            audio_path = None

    fonti_finali = filtra_fonti_per_marca(fonti_usate, filtro_marca or marca_attesa)
    if fonti_finali != sorted(fonti_usate):
        logger.info(f"Fonti filtrate per marca: {len(fonti_usate)} → {len(fonti_finali)}")

    return {
        "testo":        testo_risposta,
        "audio_path":   audio_path,
        "fonti":        fonti_finali,
        "pdf_paths":    pdf_paths,
        "pdf_richiesti": richiesta_pdf,
        "trovato":      True,
    }

# ─────────────────────────────────────────────
# NOTIFICHE ADMIN
# ─────────────────────────────────────────────
async def invia_pdf_richiesti(update: Update, context: ContextTypes.DEFAULT_TYPE, risultato: dict):
    """Invia i PDF allegati se l'utente li ha richiesti esplicitamente."""
    pdf_paths = risultato.get("pdf_paths", [])
    if not pdf_paths or not risultato.get("pdf_richiesti", False):
        return
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="upload_document")
    for pdf_path in pdf_paths[:3]:
        try:
            base = MANUALI_DIR.resolve()
            resolved = pdf_path.resolve()
            if not str(resolved).startswith(str(base)):
                logger.warning(f"Invio PDF bloccato (path non consentito): {pdf_path}")
                continue
            size_mb = pdf_path.stat().st_size / (1024 * 1024)
            if size_mb > 49:
                await update.message.reply_text(
                    f"📎 {pdf_path.name} ({size_mb:.0f}MB) — file troppo grande per Telegram."
                )
                continue
            with open(pdf_path, "rb") as pdf_file:
                await update.message.reply_document(
                    document=pdf_file,
                    filename=pdf_path.name,
                    caption=f"📄 {pdf_path.name}",
                )
            logger.info(f"PDF inviato: {pdf_path.name}")
        except Exception as e:
            logger.error(f"Errore invio PDF {pdf_path.name}: {e}")


async def notifica_admin(context, user, azione: str):
    if not ADMIN_CHAT_ID:
        return
    nome     = f"{user.first_name or ''} {user.last_name or ''}".strip()
    username = f"@{user.username}" if user.username else "nessun username"
    msg = (
        f"📚 Bot Manuali — nuova richiesta\n\n"
        f"👤 {nome} ({username})\n"
        f"🆔 ID: {user.id}\n"
        f"🕐 {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n"
        f"📌 {azione}"
    )
    try:
        await context.bot.send_message(chat_id=ADMIN_CHAT_ID, text=msg)
    except Exception as e:
        logger.error(f"Errore notifica admin: {e}")

# ─────────────────────────────────────────────
# HANDLER TELEGRAM
# ─────────────────────────────────────────────
async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    await notifica_admin(context, user, "/start")

    n_doc = collection.count() // 10 if collection.count() > 0 else 0

    await update.message.reply_text(
        f"👋 Ciao {user.first_name}!\n\n"
        f"📚 Sono il Bot Assistente Manuali Tecnici.\n\n"
        f"Posso rispondere alle tue domande tecniche sui prodotti consultando "
        f"i manuali disponibili ({n_doc} manuali indicizzati).\n\n"
        f"💬 Scrivimi direttamente la tua domanda tecnica!\n\n"
        f"📌 Comandi disponibili:\n"
        f"  /start — mostra questo messaggio\n"
        f"  /manuali — elenca i manuali disponibili\n"
        f"  /aggiorna — reindicizza i manuali nuovi o modificati (solo admin)"
    )

async def cmd_manuali(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Mostra l'elenco dei manuali indicizzati."""
    user = update.effective_user
    await notifica_admin(context, user, "/manuali")

    if collection.count() == 0:
        await update.message.reply_text("Nessun manuale ancora indicizzato.")
        return

    # Recupera i metadati a lotti per evitare "too many SQL variables"
    fonti = set()
    offset = 0
    LOTTO = 5000
    while True:
        risultati = collection.get(include=["metadatas"], limit=LOTTO, offset=offset)
        meta = risultati.get("metadatas", [])
        if not meta:
            break
        for m in meta:
            fonti.add(m.get("source", "?"))
        offset += LOTTO
        if len(meta) < LOTTO:
            break

    manuali = sorted(fonti)
    await update.message.reply_text(f"Manuali disponibili: {len(manuali)}\n\nLi invio in blocchi da 50...")

    # Invia a blocchi da 50 per non superare il limite Telegram
    blocco = ""
    count = 0
    for i, nome in enumerate(manuali, 1):
        blocco += f"{i}. {nome}\n"
        count += 1
        if count == 50:
            await update.message.reply_text(blocco)
            blocco = ""
            count = 0
    if blocco:
        await update.message.reply_text(blocco)

async def cmd_aggiorna(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Reindicizza tutti i PDF (solo per admin)."""
    user = update.effective_user

    if not ADMIN_CHAT_ID:
        await update.message.reply_text(
            "⛔ Comando disabilitato: imposta ADMIN_CHAT_ID nel file .env."
        )
        return
    if not utente_e_admin(user.id):
        await update.message.reply_text("⛔ Questo comando è riservato agli amministratori.")
        return

    await notifica_admin(context, user, "/aggiorna")
    msg = await update.message.reply_text("🔄 Reindicizzazione in corso, attendi...")

    loop = asyncio.get_event_loop()
    stats = await loop.run_in_executor(None, lambda: indicizza_manuali(forzato=False))

    await msg.edit_text(
        f"✅ Indicizzazione completata!\n\n"
        f"📂 PDF trovati: {stats['totale']}\n"
        f"🆕 Nuovi/aggiornati: {stats['nuovi']}\n"
        f"⏭️  Saltati (invariati): {stats['saltati']}\n"
        f"❌ Errori: {stats['errori']}\n\n"
        f"🔢 Chunk totali nel database: {collection.count()}"
    )

# Carica Whisper una sola volta in memoria per evitare reload ad ogni vocale
_whisper_model = None

def get_whisper_model():
    """Carica il modello Whisper una sola volta."""
    global _whisper_model
    if _whisper_model is None:
        import whisper
        logger.info("🎤 Caricamento modello Whisper medium...")
        _whisper_model = whisper.load_model("medium")
        logger.info("🎤 Modello Whisper pronto")
    return _whisper_model

async def trascrivi_audio(file_path: str) -> str:
    """Trascrive un file audio usando Whisper locale (modello medium)."""
    try:
        import asyncio
        logger.info(f"🎤 Trascrizione audio: {file_path}")

        # Prompt iniziale con termini tecnici per migliorare la trascrizione
        prompt_iniziale = (
            "Tecnico specializzato in videosorveglianza e antintrusione. "
            "Prodotti: LightSYS, WiComm, Agility, PRX, PROXIMA, SPR, PREGIO, "
            "GATEWAY2K, NVR, IPC, Dahua, RISCO, ELMO, EL.MO., "
            "centrale, sensore, rivelatore, tastiera, sirena, BUS, zona, "
            "inserimento, disinserimento, programmazione, installazione."
        )

        def _trascrivi():
            model = get_whisper_model()
            result = model.transcribe(
                file_path,
                language="it",
                initial_prompt=prompt_iniziale,
                temperature=0.0,  # deterministic, più preciso
            )
            return result["text"].strip()

        testo = await asyncio.get_event_loop().run_in_executor(None, _trascrivi)
        logger.info(f"🎤 Testo trascritto: {testo[:100]}")
        return testo
    except Exception as e:
        logger.error(f"Errore trascrizione Whisper: {e}")
        return ""

async def gestisci_vocale(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handler per messaggi vocali — trascrive con Whisper e passa a gestisci_domanda."""
    import tempfile
    user = update.effective_user

    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    msg_attesa = await update.message.reply_text(
        "🎤 Ho ricevuto il tuo messaggio vocale, trascrizione in corso..."
    )

    try:
        # Scarica il file audio da Telegram
        voice = update.message.voice or update.message.audio
        if not voice:
            await msg_attesa.edit_text("❌ Impossibile leggere il messaggio vocale.")
            return

        file = await context.bot.get_file(voice.file_id)
        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as tmp:
            tmp_path = tmp.name
        await file.download_to_drive(tmp_path)

        # Trascrivi con Whisper
        testo = await trascrivi_audio(tmp_path)

        # Elimina file temporaneo
        try:
            os.remove(tmp_path)
        except Exception:
            pass

        if not testo:
            await msg_attesa.edit_text("❌ Non sono riuscito a trascrivere il messaggio. Prova a ripetere più chiaramente.")
            return

        await msg_attesa.edit_text("🎤 Ho capito: " + testo + "\n\n🔍 Sto cercando nei manuali...")
        logger.info(f"Vocale trascritto da {user.first_name}: {testo}")

        # Notifica admin
        await notifica_admin(context, user, f"Vocale: {testo[:80]}...")

        # Aggiungi alla storia conversazione e genera risposta
        user_id = user.id
        if user_id not in CONVERSAZIONI:
            CONVERSAZIONI[user_id] = []
        CONVERSAZIONI[user_id].append({"role": "user", "content": testo})
        if len(CONVERSAZIONI[user_id]) > MAX_MESSAGGI:
            CONVERSAZIONI[user_id] = CONVERSAZIONI[user_id][-MAX_MESSAGGI:]

        risultato = await genera_risposta(testo, CONVERSAZIONI[user_id])

        if risultato["trovato"]:
            CONVERSAZIONI[user_id].append({
                "role": "assistant",
                "content": risultato["testo"][:500]
            })

        try:
            await msg_attesa.delete()
        except Exception:
            pass

        if not risultato["trovato"]:
            await update.message.reply_text(
                "❓ " + risultato["testo"] + "\n\n💡 Suggerimento: prova a includere il codice modello del prodotto."
            )
            return

        testo_msg = pulisci_testo(risultato["testo"])
        if risultato["fonti"]:
            testo_msg += "\n\nFonti:\n" + "\n".join(f"- {f}" for f in risultato["fonti"])

        if len(testo_msg) > 4000:
            for parte in [testo_msg[i:i+4000] for i in range(0, len(testo_msg), 4000)]:
                await update.message.reply_text(parte)
        else:
            await update.message.reply_text(testo_msg)

        if risultato["audio_path"] and os.path.exists(risultato["audio_path"]):
            await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="upload_voice")
            try:
                with open(risultato["audio_path"], "rb") as audio_file:
                    await update.message.reply_audio(
                        audio=audio_file,
                        caption="🔊 Risposta audio",
                        title="Risposta tecnica",
                    )
            except Exception as e:
                logger.error(f"Errore invio audio: {e}")
            finally:
                try:
                    os.remove(risultato["audio_path"])
                except Exception:
                    pass

        await invia_pdf_richiesti(update, context, risultato)

    except Exception as e:
        logger.error(f"Errore gestione vocale: {e}")
        try:
            await msg_attesa.edit_text("❌ Errore durante l'elaborazione del messaggio vocale.")
        except Exception:
            pass

async def gestisci_domanda(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handler principale: risponde a qualsiasi messaggio di testo."""
    user    = update.effective_user
    domanda = update.message.text.strip()
    if not domanda:
        return

    logger.info(f"Domanda da {user.id} ({user.username or user.first_name}): {domanda[:120]}")
    await notifica_admin(context, user, f"Domanda: {domanda[:80]}...")

    if collection.count() == 0:
        await update.message.reply_text("⚠️ Nessun manuale ancora caricato nel sistema.")
        return

    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    msg_attesa = await update.message.reply_text(
        "🔍 Sto cercando nei manuali tecnici...\n⏳ Attendi qualche secondo."
    )

    # Recupera o crea la storia della conversazione per questo utente
    user_id = user.id
    if user_id not in CONVERSAZIONI:
        CONVERSAZIONI[user_id] = []

    CONVERSAZIONI[user_id].append({"role": "user", "content": domanda})

    if len(CONVERSAZIONI[user_id]) > MAX_MESSAGGI:
        CONVERSAZIONI[user_id] = CONVERSAZIONI[user_id][-MAX_MESSAGGI:]

    risultato = await genera_risposta(domanda, CONVERSAZIONI[user_id])

    if risultato["trovato"]:
        CONVERSAZIONI[user_id].append({
            "role": "assistant",
            "content": risultato["testo"][:500]
        })

    try:
        await msg_attesa.delete()
    except Exception:
        pass

    if not risultato["trovato"]:
        await update.message.reply_text(
            f"❓ {risultato['testo']}\n\n"
            f"💡 Suggerimento: prova a includere il codice modello del prodotto."
        )
        return

    testo_msg = pulisci_testo(risultato["testo"])
    if risultato["fonti"]:
        testo_msg += "\n\nFonti:\n" + "\n".join(f"- {f}" for f in risultato["fonti"])

    if len(testo_msg) > 4000:
        for parte in [testo_msg[i:i+4000] for i in range(0, len(testo_msg), 4000)]:
            await update.message.reply_text(parte)
    else:
        await update.message.reply_text(testo_msg)

    if risultato["audio_path"] and os.path.exists(risultato["audio_path"]):
        await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="upload_voice")
        try:
            with open(risultato["audio_path"], "rb") as audio_file:
                await update.message.reply_audio(
                    audio=audio_file,
                    caption="🔊 Risposta audio",
                    title="Risposta tecnica",
                )
        except Exception as e:
            logger.error(f"Errore invio audio: {e}")
        finally:
            try:
                os.remove(risultato["audio_path"])
            except Exception:
                pass

    await invia_pdf_richiesti(update, context, risultato)

# ─────────────────────────────────────────────
# AVVIO (solo esecuzione diretta bot Telegram)
# ─────────────────────────────────────────────
if __name__ == "__main__" and _TELEGRAM_AVAILABLE:
    import atexit
    verifica_istanza_unica()
    atexit.register(rimuovi_pid_file)

    logger.info("📚 Bot Manuali Tecnici — avvio in corso...")

    # Indicizzazione in background — non blocca l'avvio del bot
    import threading
    def indicizza_in_background():
        logger.info("🔍 Controllo manuali in background...")
        stats = indicizza_manuali(forzato=False)
        logger.info(f"📊 Indicizzazione completata: {stats}")
    threading.Thread(target=indicizza_in_background, daemon=True).start()

    td_request = HTTPXRequest(connect_timeout=120, read_timeout=120)
    app = (
        ApplicationBuilder()
        .token(BOT_TOKEN)
        .request(td_request)
        .build()
    )

    app.add_handler(CommandHandler("start",    cmd_start))
    app.add_handler(CommandHandler("manuali",  cmd_manuali))
    app.add_handler(CommandHandler("aggiorna", cmd_aggiorna))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, gestisci_domanda))
    app.add_handler(MessageHandler(filters.VOICE | filters.AUDIO, gestisci_vocale))

    # Job settimanale — ogni domenica alle 03:00 scarica nuovi manuali Dahua
    import datetime as dt
    async def job_aggiornamento_settimanale(context):
        logger.info("🔄 Job settimanale: controllo nuovi manuali Dahua...")
        try:
            # Scarica le ultime pagine delle categorie principali
            categorie_principali = ["1", "100", "4572", "544", "551"]
            nuovi = 0
            for menu_id in categorie_principali:
                try:
                    r = req_lib.post(
                        DAHUA_API_SEARCH,
                        data={"id": menu_id, "keyword": "", "page": "1"},
                        headers=DAHUA_HEADERS, timeout=20
                    )
                    data = r.json().get("data", {})
                    if not isinstance(data, dict):
                        continue
                    prodotti = data.get("product", [])
                    for prod in prodotti[:5]:  # solo i più recenti
                        pid   = prod.get("product_id", "")
                        pnome = prod.get("product_name", "")
                        psafe = re.sub(r'[<>:"/\\|?*\s]', "_", pnome)[:80]
                        for tipo in ["1", "2"]:
                            nome_atteso = f"DAHUA_{psafe}_"
                            # Salta se già presente
                            esistenti = list(MANUALI_DIR.glob(f"{nome_atteso}*.pdf"))
                            if esistenti:
                                continue
                            testo, nome = scarica_pdf_dahua_online(pid, tipo, pnome)
                            if testo:
                                nuovi += 1
                            import asyncio as _aio
                            await _aio.sleep(1)
                except Exception as e:
                    logger.debug(f"Job settimanale errore categoria {menu_id}: {e}")

            logger.info(f"✅ Job settimanale completato: {nuovi} nuovi file scaricati")
            if ADMIN_CHAT_ID and nuovi > 0:
                await context.bot.send_message(
                    chat_id=ADMIN_CHAT_ID,
                    text=f"🔄 Aggiornamento settimanale completato: {nuovi} nuovi manuali scaricati."
                )
        except Exception as e:
            logger.error(f"Errore job settimanale: {e}")

    # Schedula il job ogni domenica alle 03:00
    app.job_queue.run_daily(
        job_aggiornamento_settimanale,
        time=dt.time(hour=3, minute=0),
        days=(6,),  # 6 = domenica
        name="aggiornamento_settimanale"
    )
    logger.info("⏰ Job settimanale schedulato (domenica 03:00)")

    logger.info("✅ Bot attivo e in ascolto!")
    try:
        app.run_polling(close_loop=False)
    except KeyboardInterrupt:
        logger.info("Arresto richiesto (Ctrl+C)")
    except Exception:
        logger.exception("Bot terminato per errore non gestito")
        raise SystemExit(1)
    finally:
        rimuovi_pid_file()
